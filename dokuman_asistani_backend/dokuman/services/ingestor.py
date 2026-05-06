"""Basit dosya ingest helper'larini ortak `Chunk` modeline indirger.

Bu modül; düz metin, PDF, DOCX, XLSX ve PPTX dosyalarını hızlı bir fallback
akışla okuyup açıklama/RAG katmanlarının beklediği adres, tür ve meta
alanlarına sahip chunk listeleri üretir.
"""

from __future__ import annotations

import os
import re
from dataclasses import dataclass
from typing import List, Dict, Any, Optional, Tuple

from dokuman.parcalama import chunk_text  # sende var

WORD_RE = re.compile(r"\S+")

@dataclass
class Chunk:
    """Tek bir ingest çıktısının text, adres ve metadata paketini taşır."""

    tur: str
    adres: str
    metin: str
    meta: Dict[str, Any]

def _read_text_file(path: str) -> str:
    """Metin dosyasını kontrollü encoding fallback zinciriyle okur."""
    # UTF-8 öncelikli; eski ofis çıktıları için cp1254/latin-1 fallback korunur.
    for enc in ("utf-8", "utf-8-sig", "cp1254", "latin-1"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except Exception:
            pass
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="ignore")

def _split_paragraphs(text: str) -> List[str]:
    """Ham metni boş satır sınırlarına göre anlamlı paragraf bloklarına ayırır."""
    t = (text or "").strip()
    if not t:
        return []
    # Tek satır akışlar bozulmasın diye yalnızca gerçek paragraf boşluklarında böler.
    parts = [p.strip() for p in re.split(r"\n\s*\n+", t) if p.strip()]
    return parts if parts else [t]

def _chunkify(base_adres: str, tur: str, meta: Dict[str, Any], text: str,
              chunk_char: int = 1200, overlap: int = 150) -> List[Chunk]:
    """Uzun metni adreslenebilir alt parçalara bölüp overlap bilgisini metaya yazar."""
    t = (text or "").strip()
    if not t:
        return []
    if len(t) <= int(chunk_char * 1.3):
        return [Chunk(tur=tur, adres=base_adres, metin=t, meta=meta)]
    spans = chunk_text(t, chunk_char=chunk_char, overlap=overlap)
    out: List[Chunk] = []
    for i, (b, e) in enumerate(spans, start=1):
        piece = t[b:e].strip()
        if not piece:
            continue
        out.append(Chunk(
            tur="chunk",
            adres=f"{base_adres}:chunk:{i}",
            metin=piece,
            meta={**meta, "chunk": i, "b": b, "e": e}
        ))
    return out

def ingest_txt(path: str) -> List[Chunk]:
    """TXT ve benzeri düz metinleri paragraf bazlı chunk listesine dönüştürür."""
    text = _read_text_file(path)
    paras = _split_paragraphs(text)
    out: List[Chunk] = []
    for i, p in enumerate(paras, start=1):
        out.extend(_chunkify(
            base_adres=f"txt:para:{i}",
            tur="paragraf",
            meta={"para": i},
            text=p,
        ))
    return out

def ingest_pdf(path: str) -> List[Chunk]:
    """PDF sayfalarını okuyup her sayfadaki paragrafları adresli chunk'lara çevirir."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError("PDF için pymupdf gerekli: pip install pymupdf")

    doc = fitz.open(path)
    out: List[Chunk] = []
    for page_i in range(doc.page_count):
        page = doc.load_page(page_i)
        txt = page.get_text("text") or ""
        paras = _split_paragraphs(txt)
        for pi, p in enumerate(paras, start=1):
            out.extend(_chunkify(
                base_adres=f"pdf:sayfa:{page_i+1}:para:{pi}",
                tur="paragraf",
                meta={"sayfa": page_i+1, "para": pi},
                text=p
            ))
    return out

def ingest_docx(path: str) -> List[Chunk]:
    """DOCX içindeki paragrafları ve tablo satırlarını ortak chunk sözleşmesine taşır."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("DOCX için python-docx gerekli: pip install python-docx")

    d = Document(path)
    out: List[Chunk] = []

    # Önce doğal okuma akışını veren paragraf yüzeyleri çıkarılır.
    pi = 0
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if not t:
            continue
        pi += 1
        out.extend(_chunkify(
            base_adres=f"docx:p:{pi}",
            tur="paragraf",
            meta={"p": pi},
            text=t
        ))

    # Sonra tablo satırları ayrı türle işaretlenir ki explanation katmanı formatı sezebilsin.
    ti = 0
    for table in d.tables:
        ti += 1
        for ri, row in enumerate(table.rows, start=1):
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            row_text = " | ".join([c for c in cells if c])
            if not row_text.strip():
                continue
            out.append(Chunk(
                tur="tablo_satir",
                adres=f"docx:tablo:{ti}:satir:{ri}",
                metin=row_text,
                meta={"tablo": ti, "satir": ri}
            ))

    return out

def ingest_xlsx(path: str) -> List[Chunk]:
    """Excel sayfalarını sınırlı satır/kolon taramasıyla tablo satırlarına indirger."""
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("XLSX için openpyxl gerekli: pip install openpyxl")

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out: List[Chunk] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        # Büyük dosyalarda kontrolden çıkmaması için tarama pencersi sınırlı tutulur.
        max_r = min(ws.max_row or 0, 60)
        max_c = min(ws.max_column or 0, 25)

        for r in range(1, max_r + 1):
            row_vals = []
            for c in range(1, max_c + 1):
                v = ws.cell(r, c).value
                if v is None:
                    continue
                s = str(v).strip()
                if s:
                    row_vals.append(s)
            if not row_vals:
                continue

            row_text = " | ".join(row_vals)
            out.append(Chunk(
                tur="tablo_satir",
                adres=f"xlsx:{sheet_name}:satir:{r}",
                metin=row_text,
                meta={"sheet": sheet_name, "satir": r}
            ))
    return out

def ingest_pptx(path: str) -> List[Chunk]:
    """Slayt metinlerini başlık ve madde akışını koruyan slide chunk'larına dönüştürür."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("PPTX için python-pptx gerekli: pip install python-pptx")

    prs = Presentation(path)
    out: List[Chunk] = []

    for si, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if not texts:
            continue
        slide_text = "\n".join(texts)
        out.extend(_chunkify(
            base_adres=f"pptx:slayt:{si}",
            tur="slayt",
            meta={"slayt": si},
            text=slide_text
        ))
    return out

def ingest_any(path: str, mime: str = "") -> List[Chunk]:
    """Uzantıya göre doğru ingest yolunu seçer; bilinmeyende metin fallback'i dener."""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".txt", ".md", ".py", ".js", ".json", ".csv"):
        return ingest_txt(path)
    if ext == ".pdf":
        return ingest_pdf(path)
    if ext == ".docx":
        return ingest_docx(path)
    if ext in (".xlsx", ".xlsm"):
        return ingest_xlsx(path)
    if ext == ".pptx":
        return ingest_pptx(path)

    # Bilinmeyen uzantılarda tamamen başarısız olmak yerine düz metin yolu denenir.
    return ingest_txt(path)
