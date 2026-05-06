from dataclasses import dataclass
from pathlib import Path
from typing import List, Dict, Any
import mimetypes
import re


_CODE_EXTS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cs", ".go", ".rb", ".php",
    ".sql", ".json", ".yml", ".yaml", ".html", ".css", ".xml", ".sh",
}
_TEXT_EXTS = {".txt", ".md", ".rst"}


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _slugify(text: str, *, fallback: str = "part") -> str:
    cleaned = re.sub(r"[^a-z0-9]+", "-", _clean_text(text).lower()).strip("-")
    return cleaned or fallback


def _trimmed_join(lines: list[str]) -> str:
    return "\n".join(line.rstrip() for line in lines).strip()
def _maybe_fix_mojibake(text: str) -> str:
    # BÃ–LÃœM -> BÖLÜM gibi klasik bozulma
    if any(x in text for x in ("Ã", "Ä", "Å", "Â")):
        try:
            return text.encode("latin1").decode("utf-8")
        except Exception:
            return text
    return text


def _decode_bytes_smart(data: bytes) -> str:
    # BOM kontrolü
    if data.startswith(b"\xef\xbb\xbf"):
        text = data.decode("utf-8-sig", errors="replace")
    elif data.startswith(b"\xff\xfe"):
        text = data.decode("utf-16-le", errors="replace")
    elif data.startswith(b"\xfe\xff"):
        text = data.decode("utf-16-be", errors="replace")
    else:
        # sırayla dene (latin1'i burada deneme!)
        for enc in ("utf-8", "cp1254", "iso-8859-9"):
            try:
                text = data.decode(enc)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = data.decode("utf-8", errors="replace")

    text = text.lstrip("\ufeff")
    text = text.replace("\r\n", "\n")
    return _maybe_fix_mojibake(text)
@dataclass
class Eleman:
    tur: str
    metin: str
    adres: str
    meta: Dict[str, Any]


def mime_tahmin(dosya_adi: str) -> str:
    mime, _ = mimetypes.guess_type(dosya_adi)
    return mime or ""


def txt_parcala(path: str) -> List[Eleman]:
    p = Path(path)
    data = p.read_bytes()
    text = _decode_bytes_smart(data)

    paras = [x.strip() for x in text.split("\n\n") if x.strip()]
    out = []
    for i, para in enumerate(paras, start=1):
        out.append(Eleman(tur="paragraf", metin=para, adres=f"txt:para:{i}", meta={"para": i}))
    return out


def _sheet_nonempty_rows(sheet, *, max_row: int, max_col: int):
    rows = []
    for row_idx in range(1, max_row + 1):
        values = []
        nonempty = []
        for col_idx in range(1, max_col + 1):
            raw = sheet.cell(row=row_idx, column=col_idx).value
            text = _clean_text("" if raw is None else str(raw))
            values.append(text)
            if text:
                nonempty.append(text)
        if nonempty:
            rows.append((row_idx, values, nonempty))
    return rows


def _looks_like_header_row(values: list[str]) -> bool:
    nonempty = [item for item in values if item]
    if len(nonempty) < 2:
        return False
    text_like = sum(1 for item in nonempty if re.search(r"[A-Za-zÇĞİÖŞÜçğıöşü]", item))
    numeric_like = sum(1 for item in nonempty if re.fullmatch(r"[\d.,:%+\-/]+", item))
    return text_like >= 2 and numeric_like < len(nonempty)


def pdf_parcala(path: str) -> List[Eleman]:
    import fitz  # pymupdf
    doc = fitz.open(path)
    out: List[Eleman] = []

    for pi in range(len(doc)):
        page = doc[pi]
        raw = page.get_text("text") or ""
        lines = [x.strip() for x in raw.split("\n") if x.strip()]

        # MVP birleştirme: 3 satır -> 1 parça
        para_no = 0
        buf = []
        for line in lines:
            buf.append(line)
            if len(buf) >= 3:
                para_no += 1
                out.append(Eleman(
                    tur="paragraf",
                    metin=" ".join(buf),
                    adres=f"p:{pi+1}#para:{para_no}",
                    meta={"sayfa": pi+1, "para": para_no}
                ))
                buf = []
        if buf:
            para_no += 1
            out.append(Eleman(
                tur="paragraf",
                metin=" ".join(buf),
                adres=f"p:{pi+1}#para:{para_no}",
                meta={"sayfa": pi+1, "para": para_no}
            ))

    doc.close()
    return out


def docx_parcala(path: str) -> List[Eleman]:
    from docx import Document
    d = Document(path)
    out: List[Eleman] = []
    n = 0
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if not t:
            continue
        n += 1
        style = str(p.style.name) if p.style else ""
        tur = "baslik" if style.lower().startswith("heading") else "paragraf"
        out.append(Eleman(tur=tur, metin=t, adres=f"docx:para:{n}", meta={"para": n, "stil": style}))
    return out


def xlsx_parcala(path: str) -> List[Eleman]:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    out: List[Eleman] = []
    for sheet_name in wb.sheetnames:
        sh = wb[sheet_name]
        max_row = min(sh.max_row or 0, 200)
        max_col = min(sh.max_column or 0, 30)
        rows = _sheet_nonempty_rows(sh, max_row=max_row, max_col=max_col)
        if not rows:
            continue

        header_row_idx = rows[0][0] if _looks_like_header_row(rows[0][2]) else None
        header_cells = rows[0][2][:6] if header_row_idx else []
        sheet_slug = _slugify(sheet_name, fallback="sheet")
        meta_adres = f"xlsx:sheet:{sheet_slug}#meta"
        meta_lines = [
            f"Calisma sayfasi: {sheet_name}",
            f"Toplam dolu satir: {len(rows)}",
            f"Gorulen sutun sayisi: {max_col}",
        ]
        if header_cells:
            meta_lines.append(f"Kritik sutunlar: {' | '.join(header_cells)}")
        out.append(
            Eleman(
                tur="tablo_meta",
                metin="\n".join(meta_lines),
                adres=meta_adres,
                meta={
                    "format": "xlsx",
                    "sheet": sheet_name,
                    "rows": len(rows),
                    "cols": max_col,
                    "path": meta_adres,
                    "chunk_kind": "table_meta",
                    "header_row": header_row_idx,
                    "header_cells": header_cells,
                    "baslik": sheet_name,
                },
            )
        )

        data_rows = rows[1:] if header_row_idx else rows
        if not data_rows:
            continue

        group_size = 8
        for group_index, group_start in enumerate(range(0, len(data_rows), group_size), start=1):
            group = data_rows[group_start:group_start + group_size]
            row_start = group[0][0]
            row_end = group[-1][0]
            row_lines = []
            if header_cells:
                row_lines.append(f"Basliklar: {' | '.join(header_cells)}")
            for row_idx, row_values, nonempty in group:
                visible_values = [item for item in row_values if item][:8]
                row_lines.append(f"Satir {row_idx}: {' | '.join(visible_values or nonempty[:8])}")
            adres = f"xlsx:sheet:{sheet_slug}#rows:{row_start}-{row_end}"
            out.append(
                Eleman(
                    tur="tablo",
                    metin=_trimmed_join(row_lines),
                    adres=adres,
                    meta={
                        "format": "xlsx",
                        "sheet": sheet_name,
                        "path": adres,
                        "chunk_kind": "table_rows",
                        "group_index": group_index,
                        "row_start": row_start,
                        "row_end": row_end,
                        "header_row": header_row_idx,
                        "header_cells": header_cells,
                        "baslik": f"{sheet_name} satir {row_start}-{row_end}",
                    },
                )
            )
    return out


def pptx_parcala(path: str) -> List[Eleman]:
    from pptx import Presentation
    prs = Presentation(path)
    out: List[Eleman] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        title_text = ""
        if getattr(slide.shapes, "title", None) is not None and getattr(slide.shapes.title, "text", ""):
            title_text = _clean_text(slide.shapes.title.text)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = str(shape.text or "").strip()
                if t:
                    texts.append(t)
        clean_blocks = [_clean_text(item) for item in texts if _clean_text(item)]
        if not clean_blocks:
            continue
        if not title_text:
            title_text = clean_blocks[0][:120]

        title_adres = f"pptx:slide:{i}#title"
        out.append(
            Eleman(
                tur="slayt_baslik",
                metin=title_text,
                adres=title_adres,
                meta={
                    "format": "pptx",
                    "slide": i,
                    "path": title_adres,
                    "chunk_kind": "slide_title",
                    "baslik": title_text,
                },
            )
        )

        bullet_items = []
        for item in clean_blocks:
            if item == title_text:
                continue
            bullet_items.extend([_clean_text(part) for part in re.split(r"\n+", item) if _clean_text(part)])
        if bullet_items:
            bullets_adres = f"pptx:slide:{i}#bullets"
            body_lines = [f"Baslik: {title_text}"] + [f"- {item}" for item in bullet_items[:10]]
            out.append(
                Eleman(
                    tur="slayt_madde",
                    metin="\n".join(body_lines),
                    adres=bullets_adres,
                    meta={
                        "format": "pptx",
                        "slide": i,
                        "path": bullets_adres,
                        "chunk_kind": "slide_bullets",
                        "baslik": title_text,
                        "bullet_count": len(bullet_items),
                    },
                )
            )
    return out


def _extract_symbol_name(line: str) -> str:
    patterns = [
        r"^\s*def\s+([A-Za-z_][A-Za-z0-9_]*)",
        r"^\s*class\s+([A-Za-z_][A-Za-z0-9_]*)",
        r"^\s*(?:export\s+)?(?:async\s+)?function\s+([A-Za-z_$][A-Za-z0-9_$]*)",
        r"^\s*(?:export\s+)?class\s+([A-Za-z_$][A-Za-z0-9_$]*)",
        r"^\s*(?:const|let|var)\s+([A-Za-z_$][A-Za-z0-9_$]*)\s*=\s*(?:async\s*)?\(",
        r"^\s*(SELECT|INSERT|UPDATE|DELETE|CREATE|ALTER)\b",
    ]
    for pattern in patterns:
        match = re.search(pattern, line, flags=re.IGNORECASE)
        if match:
            return str(match.group(1) or "").strip()
    return ""


def _code_block_starts(lines: list[str]) -> list[int]:
    starts = []
    for idx, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue
        if _extract_symbol_name(line):
            starts.append(idx)
    return starts


def code_parcala(path: str) -> List[Eleman]:
    p = Path(path)
    data = p.read_bytes()
    text = _decode_bytes_smart(data)
    lines = text.splitlines()
    out: List[Eleman] = []
    ext = p.suffix.lower()
    starts = _code_block_starts(lines)

    if not starts:
        starts = list(range(0, len(lines), 40))

    if starts[0] != 0:
        starts = [0] + starts

    block_no = 0
    for idx, start in enumerate(starts):
        end = starts[idx + 1] if idx + 1 < len(starts) else len(lines)
        block_lines = lines[start:end]
        block_text = _trimmed_join(block_lines)
        if not block_text:
            continue
        symbol = ""
        for line in block_lines:
            symbol = _extract_symbol_name(line)
            if symbol:
                break
        if not symbol:
            for line in block_lines:
                stripped = line.strip().lstrip("#/ -*")
                if stripped:
                    symbol = stripped[:48]
                    break
        block_no += 1
        adres = f"code:{ext.lstrip('.')}#block:{block_no}"
        out.append(
            Eleman(
                tur="kod",
                metin=block_text,
                adres=adres,
                meta={
                    "format": "code",
                    "file_ext": ext,
                    "path": adres,
                    "chunk_kind": "code_block",
                    "symbol": symbol,
                    "line_start": start + 1,
                    "line_end": end,
                    "has_comments": any(line.strip().startswith(("#", "//", "--", "/*", "*")) for line in block_lines[:6]),
                    "baslik": symbol or f"Blok {block_no}",
                },
            )
        )
    return out


def parcala(path: str) -> List[Eleman]:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return pdf_parcala(path)
    if ext == ".docx":
        return docx_parcala(path)
    if ext in [".xlsx", ".xlsm"]:
        return xlsx_parcala(path)
    if ext == ".pptx":
        return pptx_parcala(path)
    if ext in _CODE_EXTS:
        return code_parcala(path)
    if ext in _TEXT_EXTS:
        return txt_parcala(path)
    raise ValueError(f"Desteklenmeyen uzantı: {ext}")

