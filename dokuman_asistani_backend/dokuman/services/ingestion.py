"""Farkli dokuman turlerini ortak chunk ve meta sozlesmesine donusturen ingestion katmani."""

import ast
import csv
import io
import logging
import re
from collections import Counter
from pathlib import Path
from typing import Any
import xml.etree.ElementTree as ET
import zipfile

from django.conf import settings
from django.db import transaction

from ..models import Dokuman, Parca
from .parsers import mime_tahmin, _decode_bytes_smart
from dokuman.services.heading_parser import (
    looks_like_broken_ocr_line,
    normalize_section_path,
    parse_document_structure,
)
from dokuman.services.phase2_scores import analyze_cheatsheet_priority
from .ingestion_contract import (
    ingestion_bulkunu_kaydet,
    ingestion_sonucunu_kaydet,
)
from .metric_store import kaydet_skor_olayi
from .code_structure import build_code_chunk_title, build_code_segments
from dokuman.config.file_types import (
    DOCVERSE_CODE_EXTENSIONS,
    DOCVERSE_OCR_EXTENSIONS,
    DOCVERSE_PARSE_SUPPORTED_EXTENSIONS,
)

logger = logging.getLogger(__name__)

_SIGNATURE_PATTERNS = [
    r"dijital olarak imzalayan",
    r"\btarih[: ]",
    r"\+\d{2}'\d{2}'",
    r"elektronik imza",
    r"e-imza",
    r"imzalan",
    r"signature",
    r"signed by",
]
_IMAGE_EXTS = set(DOCVERSE_OCR_EXTENSIONS) & {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
_CODE_EXTS = set(DOCVERSE_CODE_EXTENSIONS) & set(DOCVERSE_PARSE_SUPPORTED_EXTENSIONS)
_TEXT_EXTS = {".txt", ".md", ".rst", ".log"}
_CSV_EXTS = {".csv", ".tsv"}
_LEGACY_OFFICE_EXTS = {".doc", ".xls", ".ppt"}
_CODE_BLOCK_RE = re.compile(
    r"^\s*(?:async\s+)?(?:def|class|function|interface|type)\s+([A-Za-z_][\w]*)"
    r"|^\s*(?:const|let|var)\s+([A-Za-z_][\w]*)\s*=\s*(?:async\s*)?\("
    r"|^\s*(?:export\s+)?(?:async\s+)?function\s+([A-Za-z_][\w]*)"
    r"|^\s*(?:public|private|protected|static|\s)*\s*([A-Za-z_][\w]*)\s*\([^)]*\)\s*\{?"
)
_COMMENT_LINE_RE = re.compile(r"^\s*(#|//|\*|/\*|\*/)")
_PPTX_PLACEHOLDER_RE = re.compile(
    r"^(?:"
    r"click to add (?:title|subtitle|text)|"
    r"click to edit master (?:title|text) style|"
    r"baslik eklemek icin tiklayin|"
    r"alt baslik eklemek icin tiklayin|"
    r"metin eklemek icin tiklayin"
    r")$",
    re.IGNORECASE,
)


def _multiformat_supported() -> set[str]:
    """PDF/DOCX disindaki parser destekli belge uzantilarini tek yerde toplar."""
    return {".xlsx", ".xlsm", ".pptx"} | set(_CODE_EXTS) | set(_TEXT_EXTS) | set(_CSV_EXTS) | set(_LEGACY_OFFICE_EXTS)


def supported_upload_extensions(*, include_images: bool = False) -> list[str]:
    """Yukleme yuzeyinde kabul edilen uzantilari tek listede toplar."""
    exts = set(getattr(settings, "DOCVERSE_UPLOAD_EXTENSIONS", [])) or ({".pdf", ".docx"} | _multiformat_supported())
    if include_images:
        exts |= _IMAGE_EXTS
    return sorted(exts)


def _legacy_office_profile(ext: str) -> dict:
    clean_ext = str(ext or "").strip().lower()
    if clean_ext == ".doc":
        return {
            "format": "doc",
            "document_type": "word",
            "unit_kind": "section",
            "conversion_target": ".docx",
        }
    if clean_ext == ".xls":
        return {
            "format": "xls",
            "document_type": "excel",
            "unit_kind": "sheet",
            "conversion_target": ".xlsx",
        }
    if clean_ext == ".ppt":
        return {
            "format": "ppt",
            "document_type": "powerpoint",
            "unit_kind": "slide",
            "conversion_target": ".pptx",
        }
    return {
        "format": clean_ext.lstrip(".") or "legacy",
        "document_type": clean_ext.lstrip(".") or "legacy",
        "unit_kind": "section",
        "conversion_target": "",
    }


def _legacy_office_error_message(ext: str, base_error: str = "") -> str:
    profile = _legacy_office_profile(ext)
    format_name = str(profile["format"]).upper()
    conversion_target = str(profile.get("conversion_target") or "").upper()
    guidance = (
        f" Legacy {format_name} dosyasi yapisal olarak sinirli okunabildi; "
        f"mumkunse {conversion_target} olarak kaydedip tekrar yukleyin."
        if conversion_target
        else f" Legacy {format_name} dosyasi yapisal olarak sinirli okunabildi."
    )
    clean_error = str(base_error or "").strip()
    if clean_error:
        clean_error = clean_error.rstrip(".")
        return f"{clean_error}.{guidance}"
    return f"Legacy {format_name} dosyasi anlamli metin uretemedi.{guidance}"


def _legacy_office_direct_support_enabled() -> bool:
    """Fragile binary Office fallback'i rollout bayragi arkasinda tutar."""
    return bool(getattr(settings, "DOCVERSE_LEGACY_OFFICE_TEXT_FALLBACK", False))


def _legacy_office_unsupported_message(ext: str) -> str:
    profile = _legacy_office_profile(ext)
    format_name = str(profile["format"]).upper()
    conversion_target = str(profile.get("conversion_target") or "").upper()
    if conversion_target:
        return (
            f"Legacy {format_name} icin dogrudan parse guvenilir degil. "
            f"Mumkunse belgeyi {conversion_target} olarak kaydedip tekrar yukleyin."
        )
    return f"Legacy {format_name} icin dogrudan parse guvenilir degil."


def _safe_ingestion_exception_text(*, ext: str = "") -> str:
    clean_ext = str(ext or "").strip().lower()
    if clean_ext in {".doc", ".xls", ".ppt"}:
        return "Legacy belge guvenli bicimde islenemedi. Mumkunse modern formata cevirip tekrar deneyin."
    return "Dokuman guvenli bicimde islenemedi."


def _redacted_exception_summary(exc) -> dict:
    return {"error_type": type(exc).__name__}


def _extract_pdf_text_layer_rows(path: str) -> list[dict]:
    """PDF'in text layer'ini sayfa bazinda inspect edip ham metni geri dondurur."""
    import fitz  # pymupdf

    rows: list[dict] = []
    doc = fitz.open(path)
    try:
        for page_no, page in enumerate(doc, start=1):
            raw_text = str(page.get_text("text") or "")
            normalized_text = _normalize_ws(raw_text)
            rows.append(
                {
                    "page": page_no,
                    "raw_text": raw_text,
                    "text": normalized_text,
                    "char_count": len(normalized_text),
                    "has_text": bool(normalized_text),
                }
            )
    finally:
        doc.close()
    return rows


def _inspect_pdf_text_layer(path: str) -> dict:
    """Scanned PDF ile text-layer PDF ayrimini hafif bir inspect adimiyla siniflandirir."""
    page_rows = _extract_pdf_text_layer_rows(path)
    page_count = len(page_rows)
    contentful_rows = [row for row in page_rows if row.get("has_text")]
    contentful_pages = len(contentful_rows)
    total_chars = sum(int(row.get("char_count") or 0) for row in contentful_rows)
    avg_chars = round(total_chars / max(contentful_pages, 1), 1) if contentful_pages else 0.0
    likely_scanned = bool(page_count > 0 and (contentful_pages == 0 or total_chars < max(90, page_count * 24)))
    return {
        "page_rows": page_rows,
        "page_count": page_count,
        "contentful_pages": contentful_pages,
        "total_chars": total_chars,
        "avg_chars_per_contentful_page": avg_chars,
        "text_layer_detected": bool(contentful_pages > 0),
        "likely_scanned": likely_scanned,
    }


def _ocr_source_profile(
    *,
    bulk: list[Parca] | None = None,
) -> dict:
    candidate_bulk = list(bulk or [])
    metas = [dict(getattr(parca, "meta", {}) or {}) for parca in candidate_bulk]
    ocr_metas = [
        meta
        for meta in metas
        if bool(meta.get("ocr_kullanildi") or meta.get("ocr") or meta.get("ocr_fallback"))
    ]
    ocr_used = bool(ocr_metas)
    kaynak_turu = ""
    if any(bool(meta.get("ocr_fallback")) for meta in ocr_metas):
        kaynak_turu = "pdf_ocr_fallback"
    elif any(bool(meta.get("ocr")) for meta in ocr_metas):
        kaynak_turu = str(
            next(
                (
                    meta.get("ocr_kaynak_turu")
                    for meta in ocr_metas
                    if str(meta.get("ocr_kaynak_turu") or "").strip()
                ),
                "image_ocr",
            )
            or "image_ocr"
        ).strip()
    elif any(bool(meta.get("text_layer_used")) for meta in metas):
        kaynak_turu = "pdf_text_layer"

    avg_ocr_quality = round(
        sum(float(meta.get("ocr_quality_score") or 0.0) for meta in ocr_metas) / max(len(ocr_metas), 1),
        3,
    ) if ocr_metas else 0.0
    if avg_ocr_quality >= 0.72:
        confidence_band = "yuksek"
    elif avg_ocr_quality >= 0.48:
        confidence_band = "orta"
    else:
        confidence_band = "dusuk" if ocr_used else ""

    warnings = [
        str(meta.get("ocr_warning") or "").strip()
        for meta in ocr_metas
        if str(meta.get("ocr_warning") or "").strip()
    ]
    warning = Counter(warnings).most_common(1)[0][0] if warnings else ""
    return {
        "ocr_kullanildi": ocr_used,
        "ocr_kaynak_turu": kaynak_turu,
        "ocr_quality_score": avg_ocr_quality,
        "ocr_confidence_band": confidence_band,
        "ocr_warning": warning,
    }


def _summarize_ocr_ingestion_signal(
    *,
    bulk: list[Parca] | None = None,
    pdf_text_layer_summary: dict | None = None,
    ocr_fallback_attempted: bool = False,
    ocr_fallback_used: bool = False,
) -> dict:
    summary = dict(pdf_text_layer_summary or {})
    profile = _ocr_source_profile(bulk=bulk)
    return {
        **profile,
        "ocr_fallback_attempted": bool(ocr_fallback_attempted),
        "ocr_fallback_used": bool(ocr_fallback_used),
        "text_layer_detected": bool(summary.get("text_layer_detected")),
        "text_layer_contentful_pages": int(summary.get("contentful_pages") or 0),
        "text_layer_total_chars": int(summary.get("total_chars") or 0),
        "likely_scanned": bool(summary.get("likely_scanned")),
    }


def _should_try_pdf_ocr_fallback(
    *,
    ext: str,
    bulk: list[Parca],
    parsed: dict,
    quality_ok: bool,
    pdf_text_layer_summary: dict | None = None,
    text_layer_bulk: list[Parca] | None = None,
) -> bool:
    if str(ext or "").lower() != ".pdf":
        return False
    if quality_ok:
        return False
    if not getattr(settings, "OCR_ENABLED", True):
        return False

    summary = dict(pdf_text_layer_summary or {})
    contentful_pages = int(summary.get("contentful_pages") or 0)
    total_layer_chars = int(summary.get("total_chars") or 0)
    candidate_bulk = list(text_layer_bulk if text_layer_bulk is not None else bulk)
    if bool(summary.get("likely_scanned")):
        return True
    if contentful_pages <= 0 or total_layer_chars < 90:
        return True

    section_count = int(parsed.get("section_count") or 0)
    contentful_count = sum(1 for parca in candidate_bulk if _has_meaningful_content(parca))
    total_chars = sum(len(_normalize_ws(parca.metin)) for parca in candidate_bulk)
    avg_quality = (
        sum(float((parca.meta or {}).get("quality_score") or 0.0) for parca in candidate_bulk) / len(candidate_bulk)
        if candidate_bulk else 0.0
    )
    weak_content_count = sum(
        1 for parca in candidate_bulk if bool((parca.meta or {}).get("weak_content"))
    )
    unique_texts = {
        _normalize_ws(getattr(parca, "metin", "")).lower()
        for parca in candidate_bulk
        if _normalize_ws(getattr(parca, "metin", ""))
    }

    if section_count <= 0 or not candidate_bulk:
        return True
    if contentful_count <= 0:
        return True
    if total_chars < 90:
        return True
    if avg_quality < 0.34:
        return True
    if candidate_bulk and (weak_content_count / max(len(candidate_bulk), 1)) >= 0.67 and avg_quality < 0.52:
        return True
    if len(candidate_bulk) >= 2 and len(unique_texts) <= 1:
        return True
    return False


def _clean_label(value: str, *, fallback: str) -> str:
    """Baslik veya etiket alanlarini tek satira indirip guvenli uzunlukta kirpar."""
    clean = re.sub(r"\s+", " ", str(value or "")).strip()
    return clean[:80] or fallback


def _truncate_words(text: str, *, limit: int = 8, fallback: str = "") -> str:
    """Uzun icerikten chunk basligi icin kisa bir kelime onizlemesi uretir."""
    words = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9_]+", str(text or ""))
    if not words:
        return fallback
    return " ".join(words[:limit])[:80]


def _language_from_ext(ext: str) -> str:
    """Dosya uzantisini code segmentation katmaninin bildigi dil etiketine cevirir."""
    mapping = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".tsx": "tsx",
        ".jsx": "jsx",
        ".java": "java",
        ".go": "go",
        ".rb": "ruby",
        ".php": "php",
        ".c": "c",
        ".cc": "cpp",
        ".cpp": "cpp",
        ".h": "c_header",
        ".hpp": "cpp_header",
        ".cs": "csharp",
        ".sql": "sql",
        ".json": "json",
        ".yml": "yaml",
        ".yaml": "yaml",
        ".html": "html",
        ".css": "css",
        ".xml": "xml",
        ".sh": "shell",
        ".ps1": "powershell",
        ".bat": "batch",
    }
    return mapping.get(ext.lower(), ext.lstrip(".") or "code")


def _build_meta_for_chunk(
    *,
    source_type: str,
    chunk_kind: str,
    adres: str,
    text: str,
    baslik: str = "",
    extra_meta: dict | None = None,
) -> dict:
    """Her chunk icin quality, difficulty ve kaynak bilgisini ayni meta semasinda birlestirir."""
    extra_meta = dict(extra_meta or {})
    office_document_type = _office_document_type_for_format(extra_meta.get("format") or source_type)
    office_unit_kind = _clean_label(
        extra_meta.get("office_unit_kind") or chunk_kind,
        fallback=chunk_kind,
    )
    quality_meta = _quality_score_analizi(
        text,
        baslik=baslik,
        icerik_uzunlugu=len(_normalize_ws(text)),
        ocr=bool(extra_meta.get("ocr")),
    )
    difficulty_meta = _difficulty_score_analizi(
        text,
        ocr=bool(extra_meta.get("ocr")),
    )
    cheatsheet_meta = analyze_cheatsheet_priority(text)
    chunk_title = _clean_label(
        baslik or extra_meta.get("chunk_title") or _truncate_words(text, fallback=adres),
        fallback=adres,
    )
    return {
        "kaynak": source_type,
        "path": adres,
        "source_address": adres,
        "format": extra_meta.get("format") or source_type,
        "document_family": f"office.{office_document_type}" if office_document_type else source_type,
        "office_document_type": office_document_type or (extra_meta.get("format") or source_type),
        "office_unit_kind": office_unit_kind,
        "office_unit_title": _clean_label(
            baslik
            or extra_meta.get("slide_title")
            or extra_meta.get("sheet")
            or extra_meta.get("chunk_title")
            or adres,
            fallback=adres,
        ),
        "chunk_kind": chunk_kind,
        "chunk_title": chunk_title,
        "baslik": baslik,
        "icerik_uzunlugu": len(_normalize_ws(text)),
        "quality_score": quality_meta["quality_score"],
        "quality_reason": quality_meta["quality_reason"],
        "weak_content": quality_meta["weak_content"],
        "quality_advisory": quality_meta["quality_advisory"],
        "short_valid": quality_meta["short_valid"],
        "difficulty_score": difficulty_meta["difficulty_score"],
        "difficulty_reason": difficulty_meta["difficulty_reason"],
        "difficulty_advisory": difficulty_meta["difficulty_advisory"],
        "cheatsheet_priority_score": cheatsheet_meta["cheatsheet_priority_score"],
        "is_cheatsheet": cheatsheet_meta["is_cheatsheet"],
        "cheatsheet_reason": cheatsheet_meta["cheatsheet_reason"],
        **extra_meta,
    }


def _office_document_type_for_format(format_name: str) -> str:
    """Format adlarini word/excel/powerpoint gibi ortak aile etiketlerine normalize eder."""
    clean = str(format_name or "").strip().lower()
    if "docx" in clean or clean == "doc" or "word" in clean:
        return "word"
    if "xlsx" in clean or "xlsm" in clean or clean == "xls" or "excel" in clean:
        return "excel"
    if "pptx" in clean or "powerpoint" in clean or clean == "ppt":
        return "powerpoint"
    return clean


def _build_multiformat_parca(
    doc: Dokuman,
    *,
    sira: int,
    tur: str,
    adres: str,
    text: str,
    source_type: str,
    chunk_kind: str,
    baslik: str = "",
    extra_meta: dict | None = None,
) -> Parca | None:
    """Ham chunk metnini kalite filtrelerinden gecirip ortak Parca nesnesine donusturur."""
    raw_text = str(text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    normalized_text = _normalize_ws(raw_text)
    if not normalized_text:
        return None
    meta = _build_meta_for_chunk(
        source_type=source_type,
        chunk_kind=chunk_kind,
        adres=adres,
        text=raw_text,
        baslik=baslik,
        extra_meta=extra_meta,
    )
    # Structured tablo/kod/liste parcalarini anlamsiz diye elemek istemiyoruz.
    if not (
        _is_meaningful_chunk(raw_text)
        or chunk_kind in {
            "table_meta",
            "table_rows",
            "table_text",
            "table_summary",
            "list_block",
            "code_comment",
            "code_block",
            "slide_title",
            "slide_bullets",
            "slide_summary",
            "slide_notes",
        }
    ):
        return None
    return Parca(
        dokuman=doc,
        sira=sira,
        tur=tur,
        metin=raw_text,
        adres=adres,
        meta=meta,
    )


def _xlsx_non_empty_rows(sheet, *, max_row: int, max_col: int) -> list[tuple[int, list[str]]]:
    """Openpyxl sheet icinden bos olmayan satirlari hucre degerleriyle birlikte toplar."""
    rows = []
    for row_idx in range(1, max_row + 1):
        values = []
        for col_idx in range(1, max_col + 1):
            cell_value = sheet.cell(row=row_idx, column=col_idx).value
            clean = _normalize_ws("" if cell_value is None else str(cell_value))
            values.append(clean)
        if any(values):
            rows.append((row_idx, values))
    return rows


def _xlsx_chunk_text(headers: list[str], rows: list[tuple[int, list[str]]]) -> str:
    """Satir gruplarini explanation katmanina okunabilir text bloklari olarak yazar."""
    header_values = [header for header in headers if header][:8]
    lines = []
    if header_values:
        lines.append(f"Basliklar: {' | '.join(header_values)}")
    for row_idx, values in rows:
        row_values = [value for value in values if value][:8]
        if not row_values:
            continue
        lines.append(f"Satir {row_idx}: {' | '.join(row_values)}")
    return "\n".join(lines)


def _is_numeric_like_value(value: str) -> bool:
    """Hucre degerinin tablo profillemede sayisal davranip davranmadigini sezgisel kontrol eder."""
    clean = _normalize_ws(value)
    if not clean:
        return False
    normalized = clean.replace("%", "").replace(",", "").replace(" ", "")
    return bool(re.fullmatch(r"[-+]?\d+(?:\.\d+)?", normalized))


def _xlsx_sheet_profile(headers: list[str], data_rows: list[tuple[int, list[str]]], *, col_count: int) -> dict:
    """Excel sheet'inin sayisal/etiket agirlikli profilini meta kararlarinda kullanmak icin cikarir."""
    profiles = []
    for col_index in range(col_count):
        header = _normalize_ws(headers[col_index] if col_index < len(headers) else "")
        values = []
        for _, row_values in data_rows:
            value = _normalize_ws(row_values[col_index] if col_index < len(row_values) else "")
            if value:
                values.append(value)
        if not header and not values:
            continue
        numeric_count = sum(1 for item in values if _is_numeric_like_value(item))
        text_count = sum(1 for item in values if not _is_numeric_like_value(item))
        profiles.append(
            {
                "header": header or f"Sutun {col_index + 1}",
                "values": values[:4],
                "numeric_count": numeric_count,
                "text_count": text_count,
                "nonempty_count": len(values),
            }
        )
    text_columns = [item["header"] for item in profiles if item["text_count"] > 0 and item["text_count"] >= item["numeric_count"]][:4]
    numeric_columns = [item["header"] for item in profiles if item["numeric_count"] > item["text_count"]][:4]
    header_value_pairs = []
    for item in profiles[:4]:
        if item["values"]:
            header_value_pairs.append(f"{item['header']}={item['values'][0]}")
    return {
        "profiles": profiles,
        "text_columns": text_columns,
        "numeric_columns": numeric_columns,
        "header_value_pairs": header_value_pairs,
    }


def _xlsx_sheet_summary_text(*, sheet_name: str, rows: list[tuple[int, list[str]]], col_count: int, header_preview: list[str], profile: dict) -> str:
    """Sheet'in ana basliklarini ve profilini explanation dostu ozet metne donusturur."""
    lines = [
        f"Calisma sayfasi {sheet_name}, {len(rows)} dolu satir ve {col_count} sutun iceriyor.",
    ]
    if header_preview:
        lines.append(f"Ana basliklar: {' | '.join(header_preview[:6])}")
    if profile.get("text_columns"):
        lines.append(f"Metin agirlikli alanlar: {' | '.join(profile['text_columns'])}")
    if profile.get("numeric_columns"):
        lines.append(f"Sayisal alanlar: {' | '.join(profile['numeric_columns'])}")
    if profile.get("header_value_pairs"):
        lines.append(f"Ornek baslik-hucre iliskileri: {' | '.join(profile['header_value_pairs'])}")
    if rows:
        first_nonempty = [value for value in rows[0][1] if value][:4]
        if first_nonempty:
            lines.append(f"Ilk anlamli satir: {' | '.join(first_nonempty)}")
    return "\n".join(lines)


def _xlsx_col_to_index(ref: str) -> int:
    """A1 benzeri hucre referansindaki kolon harflerini sifir tabanli indekse cevirir."""
    letters = "".join(ch for ch in str(ref or "") if ch.isalpha()).upper()
    idx = 0
    for ch in letters:
        idx = (idx * 26) + (ord(ch) - 64)
    return max(idx - 1, 0)


def _xlsx_shared_strings(zf: zipfile.ZipFile) -> list[str]:
    """XLSX icindeki shared string tablosunu fallback parser icin hazirlar."""
    if "xl/sharedStrings.xml" not in zf.namelist():
        return []
    root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
    ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    strings = []
    for si in root.findall("s:si", ns):
        parts = [node.text or "" for node in si.findall(".//s:t", ns)]
        strings.append(_normalize_ws("".join(parts)))
    return strings


def _xlsx_workbook_map(zf: zipfile.ZipFile) -> list[tuple[str, str]]:
    """Workbook XML'indeki sayfa adlarini ilgili sheet dosyalariyla esler."""
    ns_book = {
        "s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    ns_rel = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    workbook_root = ET.fromstring(zf.read("xl/workbook.xml"))
    rel_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
    rel_map = {
        rel.attrib.get("Id"): rel.attrib.get("Target", "")
        for rel in rel_root.findall("r:Relationship", ns_rel)
    }
    mapping = []
    for sheet in workbook_root.findall("s:sheets/s:sheet", ns_book):
        rel_id = sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        target = rel_map.get(rel_id, "")
        if not target:
            continue
        normalized_target = target.lstrip("/") if target.startswith("/") else f"xl/{target.lstrip('/')}"
        mapping.append((sheet.attrib.get("name", "Sheet"), normalized_target))
    return mapping


def _xlsx_rows_without_openpyxl(path: str) -> list[tuple[str, list[tuple[int, list[str]]], int]]:
    """Openpyxl yoksa xlsx XML'lerini okuyarak sheet satirlarini kurtarmaya calisir."""
    with zipfile.ZipFile(path) as zf:
        shared_strings = _xlsx_shared_strings(zf)
        sheet_mapping = _xlsx_workbook_map(zf)
        ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        sheets = []
        for sheet_name, target in sheet_mapping:
            if target not in zf.namelist():
                continue
            root = ET.fromstring(zf.read(target))
            rows = []
            max_col = 0
            for row in root.findall(".//s:sheetData/s:row", ns):
                row_idx = int(row.attrib.get("r") or 0)
                values_map = {}
                for cell in row.findall("s:c", ns):
                    ref = cell.attrib.get("r", "")
                    col_idx = _xlsx_col_to_index(ref)
                    cell_type = cell.attrib.get("t", "")
                    value = ""
                    raw_v = cell.find("s:v", ns)
                    inline = cell.find("s:is/s:t", ns)
                    if cell_type == "s" and raw_v is not None:
                        try:
                            value = shared_strings[int(raw_v.text or 0)]
                        except Exception:
                            value = _normalize_ws(raw_v.text or "")
                    elif inline is not None:
                        value = _normalize_ws(inline.text or "")
                    elif raw_v is not None:
                        value = _normalize_ws(raw_v.text or "")
                    if value:
                        values_map[col_idx] = value
                        max_col = max(max_col, col_idx + 1)
                if values_map:
                    row_values = [values_map.get(idx, "") for idx in range(max_col)]
                    rows.append((row_idx, row_values))
            if rows:
                sheets.append((sheet_name, rows, max_col))
        return sheets


def _pptx_slides_without_pptx(path: str) -> list[tuple[int, str, list[str]]]:
    """python-pptx yoksa zip/XML fallback'i ile slayt basligi ve maddelerini ayiklar."""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    slides = []
    with zipfile.ZipFile(path) as zf:
        slide_names = sorted(
            [name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")],
            key=_pptx_slide_sort_key,
        )
        for index, slide_name in enumerate(slide_names, start=1):
            root = ET.fromstring(zf.read(slide_name))
            texts = [_normalize_ws(node.text or "") for node in root.findall(".//a:t", ns) if _normalize_ws(node.text or "")]
            if not texts:
                continue
            slides.append((index, texts[0], texts[1:]))
    return slides


def _pptx_slide_sort_key(name: str) -> tuple[int, str]:
    """slide10/slide2 karisikliligini onlemek icin slayt adlarini sayisal siralar."""
    match = re.search(r"slide(\d+)\.xml$", str(name or ""), re.IGNORECASE)
    if match:
        return (int(match.group(1)), str(name or ""))
    return (10**9, str(name or ""))


def _is_pptx_placeholder_text(text: str) -> bool:
    """Sunum editorunun varsayilan placeholder metnini gercek icerikten ayirir."""
    clean = _normalize_ws(text)
    if not clean:
        return False
    return bool(_PPTX_PLACEHOLDER_RE.fullmatch(clean))


def _looks_like_pptx_fallback_title(text: str) -> bool:
    """Baslik placeholder'i yoksa ilk anlamli satirin baslik olup olmadigini tahmin eder."""
    clean = _normalize_ws(text)
    if not clean:
        return False

    words = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9_]+", clean)
    if not words:
        return False

    return len(words) <= 10 and len(clean) <= 80


def _resolve_pptx_slide_content(
    *,
    title_text: str,
    bullet_lines: list[str],
    other_lines: list[str],
) -> tuple[str, list[str], list[str]]:
    """Placeholder basliklari eleyip slaytin gercek baslik ve madde yapisini normalize eder."""
    title_text = "" if _is_pptx_placeholder_text(title_text) else _normalize_ws(title_text)
    bullet_lines = [
        line for line in (_normalize_ws(item) for item in (bullet_lines or []))
        if line and not _is_pptx_placeholder_text(line)
    ]
    other_lines = [
        line for line in (_normalize_ws(item) for item in (other_lines or []))
        if line and not _is_pptx_placeholder_text(line)
    ]

    if title_text:
        if bullet_lines and bullet_lines[0] == title_text:
            bullet_lines = bullet_lines[1:]
        if other_lines and other_lines[0] == title_text:
            other_lines = other_lines[1:]
        return title_text, bullet_lines, other_lines

    fallback_lines = bullet_lines or other_lines
    if len(fallback_lines) >= 2 and _looks_like_pptx_fallback_title(fallback_lines[0]):
        title_text = fallback_lines[0]
        if bullet_lines:
            bullet_lines = bullet_lines[1:]
        else:
            other_lines = other_lines[1:]

    return title_text, bullet_lines, other_lines


def _legacy_binary_candidate_lines(data: bytes) -> list[str]:
    """Binary Office dosyalarindan decode edilebilir satir adaylarini toplar."""
    candidates: list[str] = []
    seen = set()

    for match in re.finditer(rb"(?:[\x20-\x7e\xc0-\xff]\x00){6,}", data):
        try:
            text = match.group(0).decode("utf-16-le", errors="ignore")
        except Exception:
            continue
        clean = _normalize_ws(text)
        key = clean.lower()
        if clean and key not in seen:
            seen.add(key)
            candidates.append(clean)

    for match in re.finditer(rb"[\x20-\x7e\xc0-\xff\t]{8,}", data):
        try:
            text = match.group(0).decode("latin1", errors="ignore")
        except Exception:
            continue
        clean = _normalize_ws(text)
        key = clean.lower()
        if clean and key not in seen:
            seen.add(key)
            candidates.append(clean)

    filtered = []
    for item in candidates:
        if len(item) < 8 or len(item) > 240:
            continue
        if item.count("\x00") > 0:
            continue
        if re.fullmatch(r"[\W_]+", item):
            continue
        filtered.append(item)
    return filtered[:120]


def _extract_legacy_office_text(path: str) -> str:
    """Eski binary Office dosyalarindan en azindan okunabilir satirlari kurtarmaya calisir."""
    data = Path(path).read_bytes()
    lines = _legacy_binary_candidate_lines(data)
    paragraphs = []
    for line in lines:
        if _looks_like_signature_or_meta(line):
            continue
        paragraphs.append(line)
    return "\n\n".join(paragraphs[:80]).strip()


def _normalize_text_lines(text: str) -> list[str]:
    """Serbest metni heading ve blok kararlarina uygun temiz satirlara ayirir."""
    lines = []
    for raw in str(text or "").replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        lines.append(raw.rstrip())
    return lines


def _normalize_setext_headings(lines: list[str]) -> list[str]:
    """Markdown setext basliklarini ingestion icin tek tip heading satirina cevirir."""
    out: list[str] = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        next_line = lines[i + 1].strip() if i + 1 < len(lines) else ""
        if line and re.fullmatch(r"[=-]{3,}", next_line):
            out.append(f"# {line}")
            i += 2
            continue
        out.append(lines[i])
        i += 1
    return out


def _looks_like_text_heading(text: str) -> bool:
    """Kisa ve baslik benzeri satirlari paragraf bloklarindan ayirmaya yardim eder."""
    clean = _normalize_ws(text)
    if not clean:
        return False
    words = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9_]+", clean)
    if not words or len(words) > 10 or len(clean) > 90:
        return False
    if clean.startswith(("-", "*", "•")):
        return False
    if clean.endswith((".", ";")) and len(words) > 5:
        return False
    if clean.isupper():
        return True
    if re.match(r"^(#{1,6}\s+)?\d+([.)]|:)\s+", clean):
        return True
    title_ratio = sum(1 for word in words if word[:1].isupper()) / max(len(words), 1)
    return title_ratio >= 0.6


def _iter_text_blocks(text: str) -> list[dict[str, str]]:
    """Metni heading/list/code/paragraf bloklarina ayirip ingestion isaretlerini korur."""
    lines = _normalize_setext_headings(_normalize_text_lines(text))
    blocks: list[dict[str, str]] = []
    buffer: list[str] = []
    fence_lang = ""
    in_code = False

    def flush_text_buffer():
        nonlocal buffer
        chunk = "\n".join(line for line in buffer).strip()
        if chunk:
            blocks.append({"type": "text", "text": chunk})
        buffer = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code:
                code_text = "\n".join(buffer).strip()
                if code_text:
                    blocks.append({"type": "code", "text": code_text, "lang": fence_lang})
                buffer = []
                fence_lang = ""
                in_code = False
            else:
                flush_text_buffer()
                in_code = True
                fence_lang = stripped.lstrip("`").strip()
                buffer = []
            continue

        if in_code:
            buffer.append(line)
            continue

        if not stripped:
            flush_text_buffer()
            continue
        buffer.append(line)

    if in_code:
        code_text = "\n".join(buffer).strip()
        if code_text:
            blocks.append({"type": "code", "text": code_text, "lang": fence_lang})
    else:
        flush_text_buffer()
    return blocks


def _text_block_kind(lines: list[str]) -> tuple[str, str]:
    """Metin blogunu paragraf, liste veya tablo-benzeri text olarak siniflandirir."""
    nonempty = [item for item in lines if _normalize_ws(item)]
    if not nonempty:
        return "paragraf", "paragraph"
    if all(re.match(r"^\s*(?:[-*•]|\d+[.)])\s+", item) for item in nonempty):
        return "madde", "list_block"
    if len(nonempty) >= 2 and sum(1 for item in nonempty if "|" in item) >= 2:
        return "tablo", "table_text"
    return "paragraf", "paragraph"


def _build_text_bulk(
    doc: Dokuman,
    *,
    limit: int,
    source_type: str,
    format_name: str,
    text: str | None = None,
) -> tuple[list[Parca], dict]:
    """Markdown/txt benzeri kaynaklari heading-aware text ve code block chunk'larina ayirir."""
    raw_text = text if text is not None else _decode_bytes_smart(Path(doc.dosya.path).read_bytes())
    blocks = _iter_text_blocks(raw_text)
    bulk: list[Parca] = []
    sira = 1
    active_heading = ""
    heading_count = 0

    for block in blocks:
        chunk_text = str(block.get("text") or "").strip()
        if not _normalize_ws(chunk_text):
            continue
        if block.get("type") == "code":
            # Code fence bloklari textten ayri tutulur ki downstream code explanation sinyali kaybolmasin.
            language = _clean_label(block.get("lang") or format_name, fallback=format_name)
            parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur="kod",
                adres=f"{format_name}:code:{sira}",
                text=chunk_text,
                source_type=source_type,
                chunk_kind="code_block",
                baslik=active_heading or f"{language} kod blogu",
                extra_meta={
                    "format": format_name,
                    "language": language,
                    "chunk_title": active_heading or f"{language} kod blogu",
                },
            )
            if parca:
                bulk.append(parca)
                sira += 1
            if len(bulk) >= limit:
                break
            continue

        lines = [item.rstrip() for item in chunk_text.splitlines() if _normalize_ws(item)]
        if len(lines) == 1:
            heading_match = re.match(r"^\s*#{1,6}\s+(.+)$", lines[0])
            heading_text = _normalize_ws(heading_match.group(1) if heading_match else lines[0])
            if heading_match or _looks_like_text_heading(heading_text):
                # Baslik satiri kendi basina chunk olmak yerine takip eden bloga baglam tasir.
                active_heading = heading_text
                heading_count += 1
                continue

        tur, chunk_kind = _text_block_kind(lines)
        body_text = chunk_text
        if chunk_kind == "list_block":
            body_text = "\n".join(re.sub(r"^\s*(?:[-*•]|\d+[.)])\s+", "- ", item) for item in lines)
        parca = _build_multiformat_parca(
            doc,
            sira=sira,
            tur=tur,
            adres=f"{format_name}:{chunk_kind}:{sira}",
            text=body_text,
            source_type=source_type,
            chunk_kind=chunk_kind,
            baslik=active_heading,
            extra_meta={
                "format": format_name,
                "line_count": len(lines),
                "heading_title": active_heading,
                "chunk_title": active_heading or _truncate_words(body_text, fallback=f"{format_name} bolum {sira}"),
            },
        )
        if parca:
            bulk.append(parca)
            sira += 1
        if len(bulk) >= limit:
            break

    return bulk[:limit], {"block_count": len(blocks), "heading_count": heading_count}


def _read_csv_rows(path: str, *, fallback_delimiter: str = ",") -> list[list[str]]:
    """CSV/TSV icin delimiter fallback'i uygulayarak satirlari tek listede dondurur."""
    text = _decode_bytes_smart(Path(path).read_bytes())
    sample = text[:4096]
    delimiter = fallback_delimiter
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        delimiter = dialect.delimiter
    except Exception:
        delimiter = "\t" if "\t" in sample and sample.count("\t") >= sample.count(",") else fallback_delimiter
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = []
    for row in reader:
        clean = [_normalize_ws(cell) for cell in row]
        if any(clean):
            rows.append(clean)
    return rows


def _build_csv_bulk(doc: Dokuman, *, limit: int) -> tuple[list[Parca], dict]:
    """CSV dosyasini tablo ozetleri ve satir chunk'lariyla ortak bulk yapisina cevirir."""
    rows = _read_csv_rows(doc.dosya.path, fallback_delimiter="\t" if Path(doc.dosya.name).suffix.lower() == ".tsv" else ",")
    if not rows:
        return [], {"row_count": 0, "col_count": 0}

    headers = rows[0]
    data_rows = rows[1:] or rows[:1]
    header_preview = [value for value in headers if value][:6]
    col_count = max((len(row) for row in rows), default=0)
    bulk: list[Parca] = []
    sira = 1
    label = Path(doc.dosya.name).stem or "Tablo"
    meta_text = f"{label} tablosunda {len(rows)} dolu satir ve {col_count} sutun var."
    if header_preview:
        meta_text += f" Basliklar: {' | '.join(header_preview)}."
    meta_chunk = _build_multiformat_parca(
        doc,
        sira=sira,
        tur="tablo_meta",
        adres=f"csv:{label}:meta",
        text=meta_text,
        source_type="csv",
        chunk_kind="table_meta",
        baslik=label,
        extra_meta={
            "format": "csv",
            "sheet": label,
            "row_count": len(rows),
            "col_count": col_count,
            "header_preview": header_preview,
        },
    )
    if meta_chunk:
        bulk.append(meta_chunk)
        sira += 1

    for index in range(0, len(data_rows), 6):
        window = data_rows[index:index + 6]
        if not window:
            continue
        lines = []
        if header_preview:
            lines.append(f"Basliklar: {' | '.join(header_preview)}")
        for row_index, row in enumerate(window, start=index + 2):
            visible = [value for value in row if value][:8]
            lines.append(f"Satir {row_index}: {' | '.join(visible)}")
        parca = _build_multiformat_parca(
            doc,
            sira=sira,
            tur="tablo",
            adres=f"csv:{label}:rows:{index + 2}-{index + len(window) + 1}",
            text="\n".join(lines),
            source_type="csv",
            chunk_kind="table_rows",
            baslik=label,
            extra_meta={
                "format": "csv",
                "sheet": label,
                "row_start": index + 2,
                "row_end": index + len(window) + 1,
                "row_count": len(window),
                "col_count": col_count,
                "header_preview": header_preview,
            },
        )
        if parca:
            bulk.append(parca)
            sira += 1
        if len(bulk) >= limit:
            break

    return bulk[:limit], {"row_count": len(rows), "col_count": col_count}


def _build_legacy_office_bulk(doc: Dokuman, *, limit: int) -> tuple[list[Parca], dict]:
    """Legacy Office fallback'inden gelen metni list/block mantigiyla Parca bulk'una donusturur."""
    ext = Path(doc.dosya.name).suffix.lower()
    profile = _legacy_office_profile(ext)
    format_name = str(profile["format"])
    if not _legacy_office_direct_support_enabled():
        return [], {
            "legacy_format": format_name,
            "legacy_supported": False,
            "legacy_rejected": True,
            "legacy_support_mode": "conversion_required",
            "legacy_conversion_target": str(profile.get("conversion_target") or ""),
        }
    extracted_text = _extract_legacy_office_text(doc.dosya.path)
    bulk, debug_meta = _build_text_bulk(
        doc,
        limit=limit,
        source_type=f"legacy_office.{format_name}",
        format_name=format_name,
        text=extracted_text,
    )
    for index, parca in enumerate(bulk, start=1):
        meta = dict(parca.meta or {})
        source_address = str(meta.get("source_address") or parca.adres or "").strip()
        meta.update(
            {
                "format": format_name,
                "legacy_binary": True,
                "partial_parse": True,
                "structure_limited": True,
                "legacy_parse_mode": "fallback_text_extraction",
                "legacy_conversion_target": str(profile.get("conversion_target") or ""),
                "document_family": f"office.{profile['document_type']}",
                "office_document_type": str(profile["document_type"]),
                "office_unit_kind": str(profile["unit_kind"]),
                "office_unit_title": _clean_label(
                    meta.get("baslik") or meta.get("heading_title") or f"Legacy {format_name.upper()} bolum {index}",
                    fallback=f"Legacy {format_name.upper()} bolum {index}",
                ),
                "source_address": source_address or parca.adres,
                "path": source_address or parca.adres,
                "chunk_title": _clean_label(
                    meta.get("chunk_title") or meta.get("baslik") or f"Legacy {format_name.upper()} bolum {index}",
                    fallback=f"Legacy {format_name.upper()} bolum {index}",
                ),
            }
        )
        parca.meta = meta
    debug_meta["legacy_format"] = format_name
    debug_meta["legacy_binary"] = True
    debug_meta["partial_parse"] = True
    debug_meta["structure_limited"] = True
    debug_meta["legacy_conversion_target"] = str(profile.get("conversion_target") or "")
    return bulk, debug_meta


def _split_pdf_page_paragraphs(raw_text: str) -> list[str]:
    """PDF text layer'ini sayfa bazli anlamli paragraf adaylarina boler."""
    paragraphs: list[str] = []
    buffer: list[str] = []
    for raw_line in str(raw_text or "").replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        clean_line = raw_line.strip()
        if not clean_line:
            if buffer:
                paragraphs.append(" ".join(buffer).strip())
                buffer = []
            continue
        buffer.append(clean_line)
    if buffer:
        paragraphs.append(" ".join(buffer).strip())
    return [item for item in paragraphs if _normalize_ws(item)]


def _build_pdf_text_layer_bulk(
    doc: Dokuman,
    *,
    limit: int,
    page_rows: list[dict],
) -> tuple[list[Parca], dict]:
    """Heading parser zayif kaldiginda PDF text layer'ini sayfa/adres koruyarak chunk'lar."""
    bulk: list[Parca] = []
    sira = 1
    paragraph_count = 0
    contentful_pages = 0
    blocked_ocr_heading_count = 0

    for row in page_rows:
        page_no = int(row.get("page") or 0)
        paragraphs = _split_pdf_page_paragraphs(str(row.get("raw_text") or row.get("text") or ""))
        if not paragraphs and _normalize_ws(str(row.get("text") or "")):
            paragraphs = [_normalize_ws(str(row.get("text") or ""))]
        if not paragraphs:
            continue

        contentful_pages += 1
        for page_chunk_index, paragraph in enumerate(paragraphs, start=1):
            if len(bulk) >= int(limit or 0):
                break

            paragraph_count += 1
            paragraph_tokens = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9_]+", paragraph)
            ocr_heading_noise = bool(looks_like_broken_ocr_line(paragraph))
            heading_like = bool(
                _looks_like_text_heading(paragraph)
                and not ocr_heading_noise
                and len(paragraph_tokens) <= 10
                and len(paragraph) <= 90
            )
            if ocr_heading_noise and _looks_like_text_heading(paragraph):
                blocked_ocr_heading_count += 1
            tur = "baslik" if heading_like else "paragraf"
            chunk_kind = "page_heading" if heading_like else "page_paragraph"
            parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur=tur,
                adres=f"pdf:page:{page_no}:text:{page_chunk_index}",
                text=paragraph,
                source_type="pdf.text_layer",
                chunk_kind=chunk_kind,
                baslik=paragraph if heading_like else f"Sayfa {page_no}",
                extra_meta={
                    "format": "pdf",
                    "page": page_no,
                    "sayfa": page_no,
                    "page_address": f"pdf:page:{page_no}",
                    "text_layer": True,
                    "text_layer_used": True,
                    "ocr": False,
                    "ocr_fallback": False,
                    "chunk_title": paragraph if heading_like else f"PDF sayfa {page_no} paragraf {page_chunk_index}",
                },
            )
            if parca:
                bulk.append(parca)
                sira += 1
        if len(bulk) >= int(limit or 0):
            break

    return bulk[:limit], {
        "format": "pdf",
        "text_layer_used": True,
        "page_count": len(page_rows),
        "contentful_pages": contentful_pages,
        "paragraph_count": paragraph_count,
        "blocked_ocr_heading_count": blocked_ocr_heading_count,
    }


def _build_xlsx_bulk(doc: Dokuman, *, limit: int) -> tuple[list[Parca], dict]:
    """Excel dokumanlarini sheet, satir ve ozet chunk'lari halinde bulk olarak hazirlar."""
    bulk: list[Parca] = []
    sira = 1
    sheet_count = 0
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(doc.dosya.path, read_only=True, data_only=True)
        sheet_iter = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows = _xlsx_non_empty_rows(
                sheet,
                max_row=min(sheet.max_row or 0, 120),
                max_col=min(sheet.max_column or 0, 20),
            )
            sheet_iter.append((sheet_name, rows, min(sheet.max_column or 0, 20)))
    except Exception:
        sheet_iter = _xlsx_rows_without_openpyxl(doc.dosya.path)

    for sheet_name, rows, col_count in sheet_iter:
        if not rows:
            continue
        sheet_count += 1
        headers = rows[0][1]
        data_rows = rows[1:] or rows[:1]
        header_preview = [value for value in headers if value][:6]
        profile = _xlsx_sheet_profile(headers, data_rows, col_count=col_count)
        meta_text = (
            f"{sheet_name} sayfasinda {len(rows)} dolu satir ve "
            f"{col_count} sutun var."
        )
        if header_preview:
            meta_text += f" Kritik basliklar: {' | '.join(header_preview)}."
        parca = _build_multiformat_parca(
            doc,
            sira=sira,
            tur="tablo_meta",
            adres=f"xlsx:sheet:{sheet_name}:meta",
            text=meta_text,
            source_type="xlsx",
            chunk_kind="table_meta",
            baslik=f"{sheet_name} tablosu",
            extra_meta={
                "format": "xlsx",
                "sheet": sheet_name,
                "office_unit_kind": "sheet",
                "row_count": len(rows),
                "col_count": col_count,
                "header_count": len(header_preview),
                "header_preview": header_preview,
                "text_columns": list(profile.get("text_columns") or []),
                "numeric_columns": list(profile.get("numeric_columns") or []),
            },
        )
        if parca:
            bulk.append(parca)
            sira += 1
        summary_parca = _build_multiformat_parca(
            doc,
            sira=sira,
            tur="tablo_ozet",
            adres=f"xlsx:sheet:{sheet_name}:summary",
            text=_xlsx_sheet_summary_text(
                sheet_name=sheet_name,
                rows=rows,
                col_count=col_count,
                header_preview=header_preview,
                profile=profile,
            ),
            source_type="xlsx",
            chunk_kind="table_summary",
            baslik=f"{sheet_name} ozeti",
            extra_meta={
                "format": "xlsx",
                "sheet": sheet_name,
                "office_unit_kind": "sheet",
                "row_count": len(rows),
                "col_count": col_count,
                "header_preview": header_preview,
                "text_columns": list(profile.get("text_columns") or []),
                "numeric_columns": list(profile.get("numeric_columns") or []),
                "header_value_pairs": list(profile.get("header_value_pairs") or []),
            },
        )
        if summary_parca:
            bulk.append(summary_parca)
            sira += 1
        for window_start in range(0, len(data_rows), 4):
            window = data_rows[window_start:window_start + 4]
            row_start = window[0][0]
            row_end = window[-1][0]
            text = _xlsx_chunk_text(headers, window)
            parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur="tablo",
                adres=f"xlsx:sheet:{sheet_name}:rows:{row_start}-{row_end}",
                text=text,
                source_type="xlsx",
                chunk_kind="table_rows",
                baslik=f"{sheet_name} satir {row_start}-{row_end}",
                extra_meta={
                    "format": "xlsx",
                    "sheet": sheet_name,
                    "office_unit_kind": "sheet",
                    "row_start": row_start,
                    "row_end": row_end,
                    "row_count": len(window),
                    "col_count": col_count,
                    "header_preview": header_preview,
                    "text_columns": list(profile.get("text_columns") or []),
                    "numeric_columns": list(profile.get("numeric_columns") or []),
                },
            )
            if parca:
                bulk.append(parca)
                sira += 1
            if len(bulk) >= limit:
                break
        if len(bulk) >= limit:
            break
    return bulk[:limit], {"sheet_count": sheet_count}


def _pptx_shape_lines(shape) -> list[str]:
    """PowerPoint shape metnini placeholder ve bosluk gürültüsunden arindirir."""
    raw_text = str(getattr(shape, "text", "") or "")
    if not _normalize_ws(raw_text):
        return []
    return [_normalize_ws(line) for line in raw_text.splitlines() if _normalize_ws(line)]


def _build_pptx_bulk(doc: Dokuman, *, limit: int) -> tuple[list[Parca], dict]:
    """PPTX slaytlarini baslik, ozet, madde ve not yuzeylerine ayirarak chunk listesi uretir."""
    bulk: list[Parca] = []
    sira = 1
    slide_count = 0
    try:
        from pptx import Presentation

        prs = Presentation(doc.dosya.path)
        slide_iter = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            title_text = ""
            if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
                title_text = _normalize_ws(getattr(slide.shapes.title, "text", ""))
            bullet_lines: list[str] = []
            other_lines: list[str] = []
            for shape in slide.shapes:
                lines = _pptx_shape_lines(shape)
                if not lines:
                    continue
                if title_text and lines[0] == title_text:
                    continue
                # Kisa tek satirlar bullet, daha uzun akislarsa aciklama/not gibi ele alinir.
                if len(lines) == 1 and len(lines[0].split()) <= 10:
                    bullet_lines.extend(lines)
                else:
                    other_lines.extend(lines)
            title_text, bullet_lines, other_lines = _resolve_pptx_slide_content(
                title_text=title_text,
                bullet_lines=bullet_lines,
                other_lines=other_lines,
            )
            slide_iter.append((slide_index, title_text, bullet_lines, other_lines))
    except Exception:
        slide_iter = []
        for slide_index, title_text, bullet_lines in _pptx_slides_without_pptx(doc.dosya.path):
            resolved_title, resolved_bullets, resolved_other = _resolve_pptx_slide_content(
                title_text=title_text,
                bullet_lines=bullet_lines,
                other_lines=[],
            )
            slide_iter.append((slide_index, resolved_title, resolved_bullets, resolved_other))

    for slide_index, title_text, bullet_lines, other_lines in slide_iter:
        if not (title_text or bullet_lines or other_lines):
            continue
        slide_count += 1
        if title_text:
            parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur="slayt_baslik",
                adres=f"pptx:slide:{slide_index}:title",
                text=f"Slayt {slide_index} basligi: {title_text}",
                source_type="pptx",
                chunk_kind="slide_title",
                baslik=title_text,
                extra_meta={
                    "format": "pptx",
                    "slide": slide_index,
                    "office_unit_kind": "slide",
                    "slide_title": title_text,
                    "bullet_count": len(bullet_lines),
                },
            )
            if parca:
                bulk.append(parca)
                sira += 1
        summary_lines = []
        if bullet_lines:
            summary_lines.append(f"Ana maddeler: {' | '.join(bullet_lines[:3])}")
        if other_lines:
            summary_lines.append(f"Aciklama metni: {' | '.join(other_lines[:2])}")
        if summary_lines:
            summary_parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur="slayt_ozet",
                adres=f"pptx:slide:{slide_index}:summary",
                text="\n".join(
                    [f"Slayt {slide_index} ozeti: {title_text or f'Slayt {slide_index}'}"] + summary_lines
                ),
                source_type="pptx",
                chunk_kind="slide_summary",
                baslik=title_text or f"Slayt {slide_index} ozeti",
                extra_meta={
                    "format": "pptx",
                    "slide": slide_index,
                    "office_unit_kind": "slide",
                    "slide_title": title_text,
                    "bullet_count": len(bullet_lines),
                },
            )
            if summary_parca:
                bulk.append(summary_parca)
                sira += 1
        if other_lines:
            notes_parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur="slayt_not",
                adres=f"pptx:slide:{slide_index}:notes",
                text="\n".join(other_lines[:6]),
                source_type="pptx",
                chunk_kind="slide_notes",
                baslik=title_text or f"Slayt {slide_index} notlari",
                extra_meta={
                    "format": "pptx",
                    "slide": slide_index,
                    "office_unit_kind": "slide",
                    "slide_title": title_text,
                    "bullet_count": len(bullet_lines),
                },
            )
            if notes_parca:
                bulk.append(notes_parca)
                sira += 1
        content_lines = bullet_lines or other_lines
        for group_index in range(0, len(content_lines), 5):
            group = content_lines[group_index:group_index + 5]
            if not group:
                continue
            text = "\n".join(f"Madde {idx + 1}: {line}" for idx, line in enumerate(group, start=group_index))
            parca = _build_multiformat_parca(
                doc,
                sira=sira,
                tur="slayt",
                adres=f"pptx:slide:{slide_index}:bullets:{(group_index // 5) + 1}",
                text=text,
                source_type="pptx",
                chunk_kind="slide_bullets",
                baslik=title_text or f"Slayt {slide_index}",
                extra_meta={
                    "format": "pptx",
                    "slide": slide_index,
                    "office_unit_kind": "slide",
                    "slide_title": title_text,
                    "bullet_count": len(group),
                },
            )
            if parca:
                bulk.append(parca)
                sira += 1
            if len(bulk) >= limit:
                break
        if len(bulk) >= limit:
            break
    return bulk[:limit], {"slide_count": slide_count}


def _read_code_file(path: str) -> str:
    """Kod dosyasini akilli decode fallback'i ile okuyup ingestion'a hazirlar."""
    return _decode_bytes_smart(Path(path).read_bytes())


def _iter_code_segments(text: str, language: str) -> list[dict[str, Any]]:
    """Kod segmentation sorumlulugunu code_structure katmanina devreder."""
    return build_code_segments(text, language)


def _build_code_bulk(doc: Dokuman, *, limit: int) -> tuple[list[Parca], dict]:
    """Kod dosyasini diline gore segmentlere ayirip Parca nesnelerine donusturur."""
    ext = Path(doc.dosya.name).suffix.lower()
    language = _language_from_ext(ext)
    text = _read_code_file(doc.dosya.path)
    segments = _iter_code_segments(text, language)
    bulk: list[Parca] = []
    sira = 1
    unit_kind_counts: dict[str, int] = {}
    for index, segment in enumerate(segments, start=1):
        segment_text = str(segment.get("text") or "").strip()
        if not _normalize_ws(segment_text):
            continue
        unit_kind = str(segment.get("unit_kind") or "block").strip() or "block"
        unit_name = _clean_label(segment.get("unit_name") or "", fallback="")
        parent_unit = _clean_label(segment.get("parent_unit") or "", fallback="")
        symbol = unit_name or parent_unit
        unit_kind_counts[unit_kind] = unit_kind_counts.get(unit_kind, 0) + 1
        title = build_code_chunk_title(
            language,
            unit_kind,
            unit_name,
            parent_unit,
            int(segment.get("line_start") or 1),
            int(segment.get("line_end") or 1),
        )
        # Hem heuristik bayraklari hem de segmentten gelen amac ipuclarini tek metada birlestir.
        parca = _build_multiformat_parca(
            doc,
            sira=sira,
            tur="kod",
            adres=f"code:{language}:{unit_kind}:{index}",
            text=segment_text,
            source_type="code",
            chunk_kind=str(segment.get("chunk_kind") or ("code_comment" if unit_kind in {"comment", "docstring"} else "code_block")),
            baslik=title,
            extra_meta={
                "format": "code",
                "language": language,
                "code_language": language,
                "code_unit_kind": unit_kind,
                "code_unit_name": unit_name,
                "parent_unit": parent_unit,
                "code_step_kind": str(segment.get("code_step_kind") or ""),
                "test_step_kind": str(segment.get("code_step_kind") or ""),
                "code_purpose_hints": list(
                    dict.fromkeys(
                        [
                            *[
                                hint
                                for hint, enabled in (
                                    ("test", bool(segment.get("is_test"))),
                                    ("api_call", bool(segment.get("is_api_call"))),
                                    ("assertion", bool(segment.get("is_assertion"))),
                                    ("validation", bool(segment.get("is_validation"))),
                                    ("data_setup", bool(segment.get("is_data_setup"))),
                                )
                                if enabled
                            ],
                            *[str(item).strip() for item in (segment.get("purpose_hints") or []) if str(item).strip()],
                        ]
                    )
                ),
                "symbol": symbol,
                "line_start": int(segment.get("line_start") or 1),
                "line_end": int(segment.get("line_end") or 1),
                "is_test": bool(segment.get("is_test")),
                "is_api_call": bool(segment.get("is_api_call")),
                "is_assertion": bool(segment.get("is_assertion")),
                "is_validation": bool(segment.get("is_validation")),
                "is_data_setup": bool(segment.get("is_data_setup")),
                "office_unit_kind": unit_kind,
                "chunk_title": title,
            },
        )
        if parca:
            bulk.append(parca)
            sira += 1
        if len(bulk) >= limit:
            break
    return bulk[:limit], {"language": language, "segment_count": len(segments), "unit_kind_counts": unit_kind_counts}


def _validate_multiformat_bulk(bulk: list[Parca], *, source_label: str) -> tuple[bool, str]:
    """Bulk chunk'larinda bos metin ve kritik adres sorunlarini kaydetmeden once yakalar."""
    if not bulk:
        return False, f"{source_label} anlamli parca uretemedi."
    if not _paths_are_valid(bulk):
        return False, f"{source_label} gecersiz veya tekrar eden adresler uretti."
    unique_texts = {_normalize_ws(item.metin).lower() for item in bulk if _normalize_ws(item.metin)}
    if len(unique_texts) == 1 and len(bulk) >= 3:
        return False, f"{source_label} tekrar eden parcalar uretti."
    contentful = sum(1 for item in bulk if _has_meaningful_content(item))
    if contentful == 0:
        return False, f"{source_label} yalnizca zayif icerikler uretti."
    avg_quality = sum(float((item.meta or {}).get("quality_score") or 0.0) for item in bulk) / max(len(bulk), 1)
    if avg_quality < 0.28:
        return False, f"{source_label} kalite cizgisinin altinda kaldi."
    return True, ""


def _dokumani_multiformat_parcala_ve_kaydet(doc: Dokuman, *, limit: int) -> Dokuman:
    """Dosya uzantisina gore uygun parser'i secip chunklari veritabanina ve contract metriklerine yazar."""
    ext = Path(doc.dosya.name).suffix.lower()
    if ext in {".xlsx", ".xlsm"}:
        bulk, debug_meta = _build_xlsx_bulk(doc, limit=limit)
        source_type = "multiformat.xlsx"
        source_label = "XLSX parser"
        format_name = "xlsx"
    elif ext in _CSV_EXTS:
        bulk, debug_meta = _build_csv_bulk(doc, limit=limit)
        source_type = "multiformat.csv"
        source_label = "CSV parser"
        format_name = "csv"
    elif ext == ".pptx":
        bulk, debug_meta = _build_pptx_bulk(doc, limit=limit)
        source_type = "multiformat.pptx"
        source_label = "PPTX parser"
        format_name = "pptx"
    elif ext in _TEXT_EXTS:
        bulk, debug_meta = _build_text_bulk(
            doc,
            limit=limit,
            source_type=f"multiformat.{ext.lstrip('.') or 'text'}",
            format_name=ext.lstrip(".") or "text",
        )
        source_type = f"multiformat.{ext.lstrip('.') or 'text'}"
        source_label = "Metin parser"
        format_name = ext.lstrip(".") or "text"
    elif ext in _LEGACY_OFFICE_EXTS:
        bulk, debug_meta = _build_legacy_office_bulk(doc, limit=limit)
        source_type = f"multiformat.{ext.lstrip('.')}"
        source_label = "Legacy Office parser"
        format_name = ext.lstrip(".")
    elif ext in _CODE_EXTS:
        bulk, debug_meta = _build_code_bulk(doc, limit=limit)
        source_type = "multiformat.code"
        source_label = "Kod parser"
        format_name = "code"
    else:
        raise ValueError(f"Multiformat ingestion bu uzantiyi desteklemiyor: {ext}")

    quality_ok, quality_error = _validate_multiformat_bulk(bulk, source_label=source_label)
    if ext in _LEGACY_OFFICE_EXTS:
        if not _legacy_office_direct_support_enabled():
            quality_ok = False
            quality_error = _legacy_office_unsupported_message(ext)
        elif not quality_ok:
            quality_error = _legacy_office_error_message(ext, quality_error)

    persistence = ingestion_bulkunu_kaydet(
        doc,
        bulk=bulk,
        kaynak_turu=source_type,
        mime=doc.mime,
        kalite_durumu="ok" if quality_ok else "hata",
        hata_mesaji=quality_error,
        debug_ozeti={
            "format": format_name,
            "quality_ok": bool(quality_ok),
            "chunk_sayisi": len(bulk),
            "ortalama_quality_score": round(
                sum(float((p.meta or {}).get("quality_score") or 0.0) for p in bulk) / max(len(bulk), 1),
                3,
            ) if bulk else 0.0,
            "ortalama_difficulty_score": round(
                sum(float((p.meta or {}).get("difficulty_score") or 0.0) for p in bulk) / max(len(bulk), 1),
                3,
            ) if bulk else 0.0,
            **debug_meta,
        },
        batch_size=200,
        kayit_basarisiz_mesaji="Hazirlanan multiformat parcalari veritabanina kaydedilemedi.",
    )
    sonuc = persistence["sonuc"]
    real_count = persistence["real_count"]
    if getattr(doc, "owner_id", None):
        kaydet_skor_olayi(
            kullanici=doc.owner,
            olay_turu="multiformat_chunk_created",
            kaynak_modul=source_type,
            dokuman=doc,
            score_map={
                "format": format_name,
                "chunk_sayisi": len(bulk),
                "kaydedilen_parca_sayisi": real_count,
                "quality_score": round(
                    sum(float((p.meta or {}).get("quality_score") or 0.0) for p in bulk) / max(len(bulk), 1),
                    3,
                ) if bulk else 0.0,
                "difficulty_score": round(
                    sum(float((p.meta or {}).get("difficulty_score") or 0.0) for p in bulk) / max(len(bulk), 1),
                    3,
                ) if bulk else 0.0,
                "weak_content": all(bool((p.meta or {}).get("weak_content")) for p in bulk) if bulk else True,
            },
            durum="ok" if sonuc["durum_gecisi"] == "parcalandi" else "hata",
        )
    if doc.durum == "parcalandi":
        from dokuman.services.rag import sync_dokuman_indexi_if_enabled

        sync_dokuman_indexi_if_enabled(doc)
    return doc


def _normalize_ws(text: str) -> str:
    """Quality, parser ve metadata helper'lari icin tek satirlik normalize metin uretir."""
    return re.sub(r"\s+", " ", (text or "")).strip()


def _quality_score_enabled() -> bool:
    """Kalite skoru hesaplarinin feature flag ile kapatilip kapatilamayacagini soyler."""
    return bool(getattr(settings, "DOCVERSE_QUALITY_SCORE_ENABLED", True))


def _debug_summary_enabled() -> bool:
    """Debug ozeti metaya yazilsin mi sorusunu merkezi ayardan okur."""
    return bool(getattr(settings, "DOCVERSE_DEBUG_SUMMARY_ENABLED", False))


def _looks_like_signature_or_meta(text: str) -> bool:
    """Imza, tarih, tekrar eden meta veya dusuk bilgi yogunluklu satirlari eler."""
    t = _normalize_ws(text)
    if not t:
        return True

    low = t.lower()

    for pat in _SIGNATURE_PATTERNS:
        if re.search(pat, low, flags=re.IGNORECASE):
            return True

    if len(t) < 40 and re.search(r"\d{4}[./-]\d{2}[./-]\d{2}", t):
        return True

    words = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]+", t)
    if words:
        uniq_ratio = len(set(w.lower() for w in words)) / max(len(words), 1)
        if len(words) >= 6 and uniq_ratio < 0.45:
            return True

    letters = len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]", t))
    digits = len(re.findall(r"\d", t))
    if digits > 0 and letters > 0 and digits >= letters:
        return True

    return False


def _looks_like_short_meaningful_text(text: str) -> bool:
    """Kisa ama gercek icerik tasiyan satirlari kalite kapisinda korumaya calisir."""
    t = _normalize_ws(text)
    if not t:
        return False

    if _looks_like_signature_or_meta(t):
        return False

    if len(t) < 12:
        return False

    kelimeler = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9]+", t)
    harfli_kelimeler = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]{2,}", t)
    harf_sayisi = len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]", t))

    if harf_sayisi < 10:
        return False

    if len(harfli_kelimeler) < 2:
        return False

    if len(kelimeler) >= 3:
        return True

    return bool(re.search(r"[.!?;:]$", t))


def _looks_like_short_valid_content(text: str) -> bool:
    t = _normalize_ws(text)
    if not t:
        return False
    if _looks_like_signature_or_meta(t):
        return False

    if re.match(r"^[-*•]\s+[A-Za-zÇĞİÖŞÜçğıöşü0-9]", t):
        return True

    if "=" in t and len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]", t)) >= 6:
        return True

    return _looks_like_short_meaningful_text(t)


def _chunk_has_structural_value(
    text: str,
    *,
    baslik: str = "",
    meta: dict | None = None,
) -> bool:
    """Kisa ama yapisal olarak degerli parcalari yalnizca uzunluk yuzunden elememeye yardim eder."""
    t = _normalize_ws(text)
    baslik = _normalize_ws(baslik)
    meta = dict(meta or {})

    if not t or _looks_like_signature_or_meta(t):
        return False

    if _looks_like_short_valid_content(t):
        return True

    chunk_kind = str(meta.get("chunk_kind") or "").strip().lower()
    if chunk_kind in {"list_block", "table_rows", "table_summary", "code_block", "code_comment", "section"}:
        return True

    if baslik and len(t) >= 14 and len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]", t)) >= 10:
        return True

    if re.match(r"^\d+(?:[./]\d+)*\s+\S+", t):
        return True

    if "=" in t and len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü]", t)) >= 6:
        return True

    if re.search(r"[:;]", t) and len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9]+", t)) >= 3:
        return True

    return False


def _quality_score_analizi(
    text: str,
    *,
    baslik: str = "",
    icerik_uzunlugu: int = 0,
    ocr: bool = False,
) -> dict:
    t = _normalize_ws(text)
    baslik = _normalize_ws(baslik)

    if not t:
        return {
            "quality_score": 0.0,
            "quality_reason": "empty_content",
            "weak_content": True,
            "quality_advisory": True,
            "short_valid": False,
            "gate_ok": False,
            "ocr_ready": bool(ocr),
        }

    if _looks_like_signature_or_meta(t):
        return {
            "quality_score": 0.05,
            "quality_reason": "signature_or_meta",
            "weak_content": True,
            "quality_advisory": True,
            "short_valid": False,
            "gate_ok": False,
            "ocr_ready": bool(ocr),
        }

    if baslik and t.lower() == baslik.lower():
        return {
            "quality_score": 0.1,
            "quality_reason": "title_only",
            "weak_content": True,
            "quality_advisory": True,
            "short_valid": False,
            "gate_ok": False,
            "ocr_ready": bool(ocr),
        }

    words = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9]+", t)
    letters_digits = len(re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9]", t))
    alpha_num_ratio = letters_digits / max(len(t), 1)
    short_valid = (
        _looks_like_short_valid_content(t)
        or _chunk_has_structural_value(t, baslik=baslik)
        or int(icerik_uzunlugu or 0) >= 18
    )

    # Quality skoru deterministik olarak kelime sayisi ve alfa-numerik yogunluktan turetilir.
    score = min(1.0, len(words) / 20.0) * alpha_num_ratio
    if short_valid and score < 0.42:
        score = 0.42
    if ocr:
        score = min(1.0, score + 0.03)

    score = round(min(max(score, 0.0), 1.0), 3)
    weak_content = score < 0.40
    gate_ok = bool(short_valid or score >= 0.40)

    quality_reason = "contentful"
    if short_valid and len(t) < 40:
        quality_reason = "short_valid_content"
    elif len(t) < 20:
        quality_reason = "thin_content"
    elif alpha_num_ratio < 0.45:
        quality_reason = "noisy_content"

    return {
        "quality_score": score,
        "quality_reason": quality_reason,
        "weak_content": weak_content,
        "quality_advisory": weak_content or score < 0.58,
        "short_valid": short_valid,
        "gate_ok": gate_ok,
        "ocr_ready": bool(ocr),
    }


def _difficulty_score_analizi(text: str, *, ocr: bool = False) -> dict:
    t = _normalize_ws(text)
    if not t:
        return {
            "difficulty_score": 0.0,
            "difficulty_reason": "empty_content",
            "difficulty_advisory": False,
        }

    sentences = [item.strip() for item in re.split(r"[.!?;\n]+", t) if item.strip()]
    words = re.findall(r"[A-Za-zÇĞİÖŞÜçğıöşü0-9_/%+\-=]+", t)
    avg_sentence_len = len(words) / max(len(sentences), 1)
    technical_tokens = [
        token
        for token in words
        if (
            len(token) >= 10
            or re.search(r"[/%+\-=]", token)
            or re.search(r"\d", token)
            or (token.isupper() and len(token) >= 2)
        )
    ]
    technical_density = len(technical_tokens) / max(len(words), 1)
    structure_markers = len(re.findall(r"[:;()/%=]", t))
    structure_ratio = min(structure_markers / max(len(words), 1), 1.0)

    score = (
        0.45 * min(avg_sentence_len / 20.0, 1.0)
        + 0.35 * technical_density
        + 0.20 * min(structure_ratio * 3.0, 1.0)
    )
    if ocr and avg_sentence_len < 5 and technical_density < 0.15:
        score *= 0.85

    score = round(max(0.0, min(score, 1.0)), 3)
    if score >= 0.72:
        reason = "high_complexity"
    elif score >= 0.42:
        reason = "medium_complexity"
    else:
        reason = "low_complexity"

    return {
        "difficulty_score": score,
        "difficulty_reason": reason,
        "difficulty_advisory": score >= 0.58,
    }


def _is_meaningful_chunk(text: str) -> bool:
    t = _normalize_ws(text)
    if not t:
        return False
    if _looks_like_signature_or_meta(t):
        return False
    if not _quality_score_enabled():
        if len(t) < 20:
            return _looks_like_short_meaningful_text(t)
        return True
    return bool(_quality_score_analizi(t)["gate_ok"])


def _extract_content_candidate(parca: Parca) -> str:
    """Kalite dogrulamasinda kullanmak uzere parcadan metin benzeri en anlamli alani secer."""
    metin = _normalize_ws(parca.metin)
    meta = parca.meta or {}
    baslik = _normalize_ws(meta.get("baslik") or "")

    if baslik and metin.lower().startswith(baslik.lower()):
        kalan = _normalize_ws(metin[len(baslik):])
        kalan = re.sub(r"^[\s:;,\-]+", "", kalan).strip()
        return _normalize_ws(kalan)

    return metin


def _has_meaningful_content(parca: Parca) -> bool:
    icerik_uzunlugu = int((parca.meta or {}).get("icerik_uzunlugu") or 0)
    if icerik_uzunlugu >= 40:
        return True

    icerik_adayi = _extract_content_candidate(parca)
    if not icerik_adayi:
        return False

    if not _quality_score_enabled():
        if len(icerik_adayi) >= 18 and not _looks_like_signature_or_meta(icerik_adayi):
            return True
        return _looks_like_short_meaningful_text(icerik_adayi)

    analiz = _quality_score_analizi(
        icerik_adayi,
        baslik=_normalize_ws((parca.meta or {}).get("baslik") or ""),
        icerik_uzunlugu=icerik_uzunlugu,
    )
    return bool(
        analiz["gate_ok"]
        or _chunk_has_structural_value(
            icerik_adayi,
            baslik=_normalize_ws((parca.meta or {}).get("baslik") or ""),
            meta=parca.meta,
        )
    )


def _normalize_address_for_validation(adres: str) -> str:
    clean = _normalize_ws(adres)
    if not clean:
        return ""
    if re.fullmatch(r"[0-9.\s]+", clean):
        return normalize_section_path(clean, fallback="")
    return clean


def _paths_are_valid(bulk: list[Parca]) -> bool:
    adresler = [
        _normalize_address_for_validation(getattr(parca, "adres", ""))
        for parca in bulk
    ]
    if not adresler or any(not adres for adres in adresler):
        return False
    return len(adresler) == len(set(adresler))


def _validate_bulk_quality(bulk: list[Parca], parsed: dict, limit: int) -> tuple[bool, str]:
    section_count = int(parsed.get("section_count") or 0)
    if not bulk:
        if section_count > 0:
            return False, "Heading parser bölüm çıkardı ama ingestion kalite filtresi tüm parçaları eledi."
        return False, "Heading parser anlamlı parça üretemedi."

    if section_count <= 0:
        return False, "Heading parser bölüm yapısı çıkaramadı."

    limited_bulk = bulk[:limit]
    total_chars = sum(len(_normalize_ws(p.metin)) for p in limited_bulk)
    contentful_count = sum(
        1
        for p in limited_bulk
        if _has_meaningful_content(p)
    )
    structural_count = sum(
        1
        for p in limited_bulk
        if _chunk_has_structural_value(
            _extract_content_candidate(p),
            baslik=_normalize_ws((p.meta or {}).get("baslik") or ""),
            meta=p.meta,
        )
    )
    quality_scores = [
        float(((p.meta or {}).get("quality_score")) or _quality_score_analizi(
            _extract_content_candidate(p),
            baslik=_normalize_ws((p.meta or {}).get("baslik") or ""),
            icerik_uzunlugu=int((p.meta or {}).get("icerik_uzunlugu") or 0),
            ocr=bool((p.meta or {}).get("ocr")),
        )["quality_score"])
        for p in limited_bulk
    ]
    weak_content_count = sum(
        1 for p in limited_bulk if bool((p.meta or {}).get("weak_content"))
    )
    unique_texts = {
        _normalize_ws(p.metin).lower()
        for p in limited_bulk
        if _normalize_ws(p.metin)
    }

    min_total_chars = 48 if contentful_count >= 2 else 80
    if total_chars < min_total_chars and structural_count == 0:
        return False, "Heading parser çok az metin çıkardı."

    if contentful_count == 0 and structural_count == 0:
        return False, "Heading parser yalnız başlık benzeri parçalar çıkardı."

    if quality_scores and max(quality_scores) < 0.42:
        return False, "Heading parser çok zayıf içerikler çıkardı."

    if len(limited_bulk) >= 3 and len(unique_texts) <= 1:
        return False, "Heading parser tekrar eden parçalar üretti."

    if not _paths_are_valid(limited_bulk):
        return False, "Heading parser geçersiz veya tekrar eden adresler üretti."

    if weak_content_count == len(limited_bulk) and quality_scores and sum(quality_scores) / len(quality_scores) < 0.55:
        return False, "Heading parser zayif icerikleri agirlikli olarak çıkardı."

    return True, ""


def _build_bulk_from_heading_parser(doc: Dokuman):
    """Heading parser bolumlerini kalite metasi korunmus Parca listesine cevirir."""
    ext = Path(doc.dosya.name).suffix.lower()
    format_name = ext.lstrip(".") or "document"
    document_type = _office_document_type_for_format(format_name)
    document_family = f"office.{document_type}" if document_type in {"word", "excel", "powerpoint"} else document_type

    if ext not in {".pdf", ".docx"}:
        raise ValueError(f"Heading parser bu uzantıyı desteklemiyor: {ext}")

    parsed = parse_document_structure(doc.dosya.path)
    sections = parsed.get("sections") or []

    bulk = []
    sira = 1

    for sec in sections:
        title = (sec.get("title") or "").strip()
        content = (sec.get("content") or "").strip()
        level = sec.get("level")
        page_start = sec.get("page_start")
        path_value = normalize_section_path((sec.get("path") or "").strip(), fallback=str(sira))

        if title and content:
            parca_metin = f"{title}\n\n{content}".strip()
        else:
            parca_metin = (title or content).strip()

        parca_metin = _normalize_ws(parca_metin)
        quality_meta = _quality_score_analizi(
            content or parca_metin,
            baslik=title,
            icerik_uzunlugu=len(_normalize_ws(content)),
            ocr=False,
        )
        difficulty_meta = _difficulty_score_analizi(content or parca_metin)
        cheatsheet_meta = analyze_cheatsheet_priority(content or parca_metin)

        if not (
            bool(quality_meta["gate_ok"])
            or _chunk_has_structural_value(content or parca_metin, baslik=title, meta={"chunk_kind": "section"})
        ):
            continue

        safe_address = path_value or f"section:{sira}"
        bulk.append(
            Parca(
                dokuman=doc,
                sira=sira,
                tur="bolum",
                metin=parca_metin,
                adres=safe_address,
                meta={
                    "kaynak": "heading_parser",
                    "format": format_name,
                    "document_family": document_family,
                    "office_document_type": document_type,
                    "office_unit_kind": "section",
                    "office_unit_title": title or safe_address or f"Bolum {sira}",
                    "chunk_kind": "section",
                    "chunk_title": title or safe_address or f"Bolum {sira}",
                    "source_address": safe_address,
                    "baslik": title,
                    "seviye": level,
                    "sayfa": page_start,
                    "path": safe_address,
                    "icerik_uzunlugu": len(_normalize_ws(content)),
                    "quality_score": quality_meta["quality_score"],
                    "quality_reason": quality_meta["quality_reason"],
                    "weak_content": quality_meta["weak_content"],
                    "quality_advisory": quality_meta["quality_advisory"],
                    "short_valid": quality_meta["short_valid"],
                    "difficulty_score": difficulty_meta["difficulty_score"],
                    "difficulty_reason": difficulty_meta["difficulty_reason"],
                    "difficulty_advisory": difficulty_meta["difficulty_advisory"],
                    # Faz 2 icin hap bilgiye uygun parcalari statik olarak isaretliyoruz.
                    "cheatsheet_priority_score": cheatsheet_meta["cheatsheet_priority_score"],
                    "is_cheatsheet": cheatsheet_meta["is_cheatsheet"],
                    "cheatsheet_reason": cheatsheet_meta["cheatsheet_reason"],
                    "heading_score": sec.get("heading_score"),
                    "heading_decision_reason": sec.get("heading_decision_reason") or "",
                },
            )
        )
        sira += 1

    return bulk, parsed

def dokumani_parcala_ve_kaydet(doc: Dokuman, limit: int = 5000) -> Dokuman:
    """
    Dosya turune gore uygun ingestion hattini secip chunk, meta ve kalite kaydini tamamlar.
    """
    doc.mime = doc.mime or mime_tahmin(doc.dosya.name)
    ext = Path(doc.dosya.name).suffix.lower()
    pdf_text_layer_summary: dict = {}
    text_layer_bulk: list[Parca] = []

    if ext in _IMAGE_EXTS:
        from .ocr import gorseli_ocr_ile_parcala_ve_kaydet

        return gorseli_ocr_ile_parcala_ve_kaydet(doc)

    if ext in _multiformat_supported():
        return _dokumani_multiformat_parcala_ve_kaydet(doc, limit=limit)

    try:
        if ext == ".pdf":
            try:
                pdf_text_layer_summary = _inspect_pdf_text_layer(doc.dosya.path)
            except Exception as inspect_error:
                pdf_text_layer_summary = {
                    "page_rows": [],
                    "page_count": 0,
                    "contentful_pages": 0,
                    "total_chars": 0,
                    "avg_chars_per_contentful_page": 0.0,
                    "text_layer_detected": False,
                    "likely_scanned": True,
                    "inspect_error_type": type(inspect_error).__name__,
                }
        logger.debug("Ingestion heading parser start for doc_id=%s ext=%s", getattr(doc, "id", None), ext)

        bulk, parsed = _build_bulk_from_heading_parser(doc)

        if len(bulk) > limit:
            bulk = bulk[:limit]

        logger.debug(
            "Ingestion heading parser done for doc_id=%s ext=%s bulk=%s section_count=%s",
            getattr(doc, "id", None),
            ext,
            len(bulk),
            parsed.get("section_count", 0),
        )

        quality_ok, quality_error = _validate_bulk_quality(bulk, parsed, limit)
        logger.debug("Ingestion quality check for doc_id=%s ext=%s quality_ok=%s", getattr(doc, "id", None), ext, quality_ok)
        if ext == ".pdf" and not quality_ok and pdf_text_layer_summary.get("text_layer_detected"):
            text_layer_bulk, text_layer_debug = _build_pdf_text_layer_bulk(
                doc,
                limit=limit,
                page_rows=list(pdf_text_layer_summary.get("page_rows") or []),
            )
            text_layer_ok, text_layer_error = _validate_multiformat_bulk(
                text_layer_bulk,
                source_label="PDF text layer parser",
            )
            if text_layer_ok:
                persistence = ingestion_bulkunu_kaydet(
                    doc,
                    bulk=text_layer_bulk,
                    kaynak_turu="pdf.text_layer",
                    mime=doc.mime,
                    kalite_durumu="ok",
                    hata_mesaji="",
                    debug_ozeti={
                        "section_count": int(parsed.get("section_count") or 0),
                        "heading_parser_quality_ok": bool(quality_ok),
                        "heading_parser_hata": quality_error,
                        **_summarize_ocr_ingestion_signal(
                            bulk=text_layer_bulk,
                            pdf_text_layer_summary=pdf_text_layer_summary,
                            ocr_fallback_attempted=False,
                            ocr_fallback_used=False,
                        ),
                        **text_layer_debug,
                    },
                    batch_size=500,
                    kayit_basarisiz_mesaji="PDF text layer parcalari veritabanina kaydedilemedi.",
                )
                sonuc = persistence["sonuc"]
                real_count = persistence["real_count"]
                if getattr(doc, "owner_id", None):
                    kaydet_skor_olayi(
                        kullanici=doc.owner,
                        olay_turu="multiformat_chunk_created",
                        kaynak_modul="pdf.text_layer",
                        dokuman=doc,
                        score_map={
                            "format": "pdf",
                            "chunk_sayisi": len(text_layer_bulk),
                            "kaydedilen_parca_sayisi": real_count,
                            "quality_score": round(
                                sum(float((p.meta or {}).get("quality_score") or 0.0) for p in text_layer_bulk)
                                / max(len(text_layer_bulk), 1),
                                3,
                            ) if text_layer_bulk else 0.0,
                            "difficulty_score": round(
                                sum(float((p.meta or {}).get("difficulty_score") or 0.0) for p in text_layer_bulk)
                                / max(len(text_layer_bulk), 1),
                                3,
                            ) if text_layer_bulk else 0.0,
                            "weak_content": all(
                                bool((p.meta or {}).get("weak_content")) for p in text_layer_bulk
                            ) if text_layer_bulk else True,
                            "text_layer_used": True,
                            "ocr_fallback": False,
                        },
                        durum="ok" if sonuc["durum_gecisi"] == "parcalandi" else "hata",
                    )
                if doc.durum == "parcalandi":
                    from dokuman.services.rag import sync_dokuman_indexi_if_enabled

                    sync_dokuman_indexi_if_enabled(doc)
                return doc
            quality_error = text_layer_error or quality_error

        if _should_try_pdf_ocr_fallback(
            ext=ext,
            bulk=bulk,
            parsed=parsed,
            quality_ok=quality_ok,
            pdf_text_layer_summary=pdf_text_layer_summary,
            text_layer_bulk=text_layer_bulk,
        ):
            from .ocr import pdf_ocr_ile_parcala_ve_kaydet

            return pdf_ocr_ile_parcala_ve_kaydet(
                doc,
                limit=limit,
                fallback_reason=quality_error,
                pdf_text_layer_summary=pdf_text_layer_summary,
            )

        parca_sayisi = len(bulk)
        logger.debug("Ingestion bulk ready for doc_id=%s ext=%s bulk=%s", getattr(doc, "id", None), ext, parca_sayisi)

        persistence = ingestion_bulkunu_kaydet(
            doc,
            bulk=bulk,
            kaynak_turu="heading_parser",
            mime=doc.mime,
            kalite_durumu="ok" if quality_ok else "hata",
            hata_mesaji=quality_error,
            debug_ozeti={
                "section_count": int(parsed.get("section_count") or 0),
                "quality_ok": bool(quality_ok),
                **_summarize_ocr_ingestion_signal(
                    bulk=bulk,
                    pdf_text_layer_summary=pdf_text_layer_summary,
                    ocr_fallback_attempted=(ext == ".pdf" and not quality_ok),
                    ocr_fallback_used=False,
                ),
                **(
                    {
                        "ortalama_quality_score": round(
                            sum(float((p.meta or {}).get("quality_score") or 0.0) for p in bulk) / max(len(bulk), 1),
                            3,
                        ) if bulk else 0.0,
                        "ortalama_difficulty_score": round(
                            sum(float((p.meta or {}).get("difficulty_score") or 0.0) for p in bulk) / max(len(bulk), 1),
                            3,
                        ) if bulk else 0.0,
                        "weak_content_parca_sayisi": sum(
                            1 for p in bulk if bool((p.meta or {}).get("weak_content"))
                        ),
                        "cheatsheet_parca_sayisi": sum(
                            1 for p in bulk if bool((p.meta or {}).get("is_cheatsheet"))
                        ),
                        "heading_reason_ozeti": dict(parsed.get("debug_ozeti", {}).get("heading_reason_ozeti") or {}),
                    }
                    if _debug_summary_enabled()
                    else {}
                ),
            },
            batch_size=500,
            kayit_basarisiz_mesaji="Hazirlanan parçalar veritabanina kaydedilemedi.",
        )
        sonuc = persistence["sonuc"]
        real_count = persistence["real_count"]
        logger.debug(
            "Ingestion persisted for doc_id=%s ext=%s saved=%s state=%s",
            getattr(doc, "id", None),
            ext,
            real_count,
            sonuc["durum_gecisi"],
        )

        if getattr(doc, "owner_id", None):
            kaydet_skor_olayi(
                kullanici=doc.owner,
                olay_turu="ingestion_skorlandi",
                kaynak_modul="ingestion.heading_parser",
                dokuman=doc,
                score_map={
                    "quality_score": (
                        sum(float((p.meta or {}).get("quality_score") or 0.0) for p in bulk) / max(len(bulk), 1)
                    ) if bulk else 0.0,
                    "difficulty_score": (
                        sum(float((p.meta or {}).get("difficulty_score") or 0.0) for p in bulk) / max(len(bulk), 1)
                    ) if bulk else 0.0,
                    "quality_reason": "contentful" if quality_ok else "quality_gate_failed",
                    "difficulty_reason": "aggregate_chunk_difficulty",
                    "kaydedilen_parca_sayisi": real_count,
                },
                durum="ok" if doc.durum == "parcalandi" else "hata",
            )

        if doc.durum == "parcalandi":
            from dokuman.services.rag import sync_dokuman_indexi_if_enabled

            sync_dokuman_indexi_if_enabled(doc)

        return doc

    except Exception as e:
        ingestion_sonucunu_kaydet(
            doc,
            kaynak_turu="heading_parser",
            mime=doc.mime,
            aday_parca_sayisi=0,
            kaydedilen_parca_sayisi=0,
            kalite_durumu="hata",
            hata_mesaji=_safe_ingestion_exception_text(ext=ext),
            debug_ozeti=_redacted_exception_summary(e),
        )
        logger.warning(
            "Ingestion failed for doc_id=%s ext=%s error_type=%s",
            getattr(doc, "id", None),
            ext,
            type(e).__name__,
        )
        raise
